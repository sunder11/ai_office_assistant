# Add these imports with the existing ones
from pptx import Presentation
import pandas as pd
import openpyxl
import os
from dotenv import load_dotenv
import streamlit as st
from docling.document_converter import DocumentConverter
from langchain.schema import Document
from langchain_text_splitters.character import CharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_groq import ChatGroq
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
import xml.etree.ElementTree as ET
from typing import List
from urllib.parse import urljoin
import requests
import time
import random
import sqlite3
import datetime
import threading
from pathlib import Path
import glob
import re

import time
import pickle
import os
from pathlib import Path


# Load the environment variables
load_dotenv()
working_dir = os.path.dirname(os.path.abspath(__file__))
# Define data directory for persistent storage
DATA_DIR = os.path.join(working_dir, "data")
os.makedirs(DATA_DIR, exist_ok=True)
DB_PATH = os.path.join(DATA_DIR, "website_data.db")
# Thread-local storage for database connections
_local = threading.local()


### LOGIN
# Simple file-based session storage
SESSION_FILE = os.path.join(DATA_DIR, "session.pkl")


def check_login_credentials(username: str, password: str) -> bool:
    """Check if login credentials match environment variables"""
    env_username = os.getenv("USER_USER")
    env_password = os.getenv("USER_PASS")

    if not env_username or not env_password:
        st.error("Login credentials not configured in environment variables")
        return False

    return username == env_username and password == env_password


def validate_login_input(username: str, password: str) -> tuple[bool, str]:
    """Validate login input format"""
    if not username or not password:
        return False, "Please enter both username and password"

    if len(username) > 10:
        return False, "Username must be 10 characters or less"

    if len(password) < 9:
        return False, "Password must be at least 9 characters"

    if len(password) > 25:
        return False, "Password must be 25 characters or less"

    return True, ""


def save_session_to_file(username: str):
    """Save session to file"""
    try:
        session_data = {
            "username": username,
            "login_time": time.time(),
            "authenticated": True,
        }
        with open(SESSION_FILE, "wb") as f:
            pickle.dump(session_data, f)
        return True
    except Exception as e:
        st.error(f"Error saving session: {e}")
        return False


def load_session_from_file():
    """Load session from file"""
    try:
        if not os.path.exists(SESSION_FILE):
            return None

        with open(SESSION_FILE, "rb") as f:
            session_data = pickle.load(f)

        # Check if session is still valid (24 hours)
        login_time = session_data.get("login_time", 0)
        hours_elapsed = (time.time() - login_time) / 3600

        if hours_elapsed < 24 and session_data.get("authenticated", False):
            return session_data
        else:
            # Session expired, remove file
            clear_session_file()
            return None
    except Exception as e:
        # File corrupted or error, remove it
        clear_session_file()
        return None


def clear_session_file():
    """Clear session file"""
    try:
        if os.path.exists(SESSION_FILE):
            os.remove(SESSION_FILE)
    except Exception:
        pass


def create_persistent_login(username: str):
    """Create persistent login session"""
    # Set Streamlit session state
    st.session_state.authenticated = True
    st.session_state.username = username
    st.session_state.login_timestamp = time.time()
    st.session_state.session_valid = True

    # Save to file for persistence
    save_session_to_file(username)


def check_persistent_login():
    """Check if there's a valid persistent login"""
    # First check if already authenticated in current session
    if st.session_state.get("authenticated", False) and st.session_state.get(
        "session_valid", False
    ):
        # Double-check the session is still valid
        login_time = st.session_state.get("login_timestamp", 0)
        hours_elapsed = (time.time() - login_time) / 3600
        if hours_elapsed < 24:
            return True
        else:
            logout_silent()

    # If not authenticated in session state, try to load from file
    session_data = load_session_from_file()
    if session_data:
        # Restore session
        st.session_state.authenticated = True
        st.session_state.username = session_data["username"]
        st.session_state.login_timestamp = session_data["login_time"]
        st.session_state.session_valid = True
        return True

    return False


def logout_silent():
    """Clear authentication without rerun"""
    keys_to_clear = ["authenticated", "username", "login_timestamp", "session_valid"]
    for key in keys_to_clear:
        if key in st.session_state:
            del st.session_state[key]

    # Clear session file
    clear_session_file()


def show_login_form():
    """Display login form and handle authentication"""
    st.set_page_config(
        page_title="Login - Claudia User", page_icon="🔒", layout="centered"
    )

    # Apply CSS to hide menu and deploy button on login page
    st.markdown(
        """
        <style>
            .reportview-container {
                margin-top: -2em;
            }
            #MainMenu {visibility: hidden;}
            .stDeployButton {display:none;}
            footer {visibility: hidden;}
            #stDecoration {display:none;}
            
            /* Login form styling */
            .stForm {
                background-color: #f8f9fa;
                padding: 2rem;
                border-radius: 10px;
                border: 2px solid #dee2e6;
                box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
            }
        </style>
    """,
        unsafe_allow_html=True,
    )

    # Center the login form
    col1, col2, col3 = st.columns([1, 2, 1])

    with col2:
        st.title("🔒 Login Required")
        st.markdown("---")

        # Create login form
        with st.form("login_form", clear_on_submit=False):
            st.subheader("Claudia User Access")

            username = st.text_input(
                "Username",
                max_chars=10,
                placeholder="Enter username (max 10 chars)",
                help="Maximum 10 characters",
            )

            password = st.text_input(
                "Password",
                type="password",
                max_chars=25,
                placeholder="Enter password (9-25 chars)",
                help="Minimum 9 characters, maximum 25 characters",
            )

            login_button = st.form_submit_button(
                "Login", type="primary", use_container_width=True
            )

            if login_button:
                # Validate input format
                is_valid, error_message = validate_login_input(username, password)

                if not is_valid:
                    st.error(error_message)
                    return False

                # Check credentials
                if check_login_credentials(username, password):
                    create_persistent_login(username)
                    st.success("Login successful! Redirecting...")
                    time.sleep(0.5)
                    st.rerun()
                else:
                    st.error("Invalid username or password")
                    return False

    return False


def check_authentication():
    """Check if user is authenticated, show login if not"""
    # Check for persistent login first
    if check_persistent_login():
        return True

    # Show login form if not authenticated
    show_login_form()
    return False


def logout():
    """Clear authentication and force re-login"""
    logout_silent()
    st.rerun()


def get_all_names_from_vectorstore_improved(vectorstore):
    """Extract ONLY real names from the vectorstore - FIXED for real names with punctuation"""
    try:
        all_names = set()

        if hasattr(vectorstore, "docstore") and hasattr(
            vectorstore, "index_to_docstore_id"
        ):
            for i in range(len(vectorstore.index_to_docstore_id)):
                doc_id = vectorstore.index_to_docstore_id.get(i)
                if doc_id and doc_id in vectorstore.docstore._dict:
                    doc = vectorstore.docstore._dict[doc_id]

                    # Only process Excel chunks
                    if doc.metadata.get("file_type") in ["xlsx", "xls"]:
                        content = doc.page_content

                        # ONLY use the most precise pattern for your data structure
                        name_pattern = (
                            r"FIRST_NAME:\s*([^|]+?)\s*\|\s*LAST_NAME:\s*([^|]+?)\s*\|"
                        )
                        matches = re.findall(name_pattern, content, re.IGNORECASE)

                        for first, last in matches:
                            first = first.strip()
                            last = last.strip()

                            # More realistic filtering for real names
                            if (
                                first
                                and last
                                and len(first) >= 2
                                and len(last) >= 2
                                and
                                # Allow letters, spaces, periods, commas, apostrophes, hyphens
                                re.match(r"^[A-Za-z\s\.\,\'\-]+$", first)
                                and re.match(
                                    r"^[A-Za-z\s\.\,\'\-IVX]+$", last
                                )  # Allow Roman numerals
                                and
                                # Exclude obvious non-names
                                not any(
                                    word.lower() in first.lower()
                                    for word in [
                                        "record",
                                        "structured",
                                        "data",
                                        "sheet",
                                        "column",
                                        "total",
                                        "fee",
                                    ]
                                )
                                and not any(
                                    word.lower() in last.lower()
                                    for word in [
                                        "record",
                                        "structured",
                                        "data",
                                        "sheet",
                                        "column",
                                        "total",
                                        "fee",
                                    ]
                                )
                                and
                                # Must start with a letter
                                first[0].isalpha()
                                and last[0].isalpha()
                            ):
                                full_name = f"{first} {last}"
                                all_names.add(full_name)

        return sorted(list(all_names))

    except Exception as e:
        st.error(f"Error extracting names: {str(e)}")
        return []


def get_db_connection():
    """Get a thread-specific database connection"""
    if not hasattr(_local, "conn"):
        _local.conn = sqlite3.connect(DB_PATH, check_same_thread=False)
        # Enable foreign keys
        _local.conn.execute("PRAGMA foreign_keys = ON")
    return _local.conn


def close_db_connection():
    """Close the thread-specific database connection if it exists"""
    if hasattr(_local, "conn"):
        _local.conn.close()
        delattr(_local, "conn")


def setup_database():
    """Set up SQLite database for document storage"""
    conn = get_db_connection()
    cursor = conn.cursor()
    # Create enhanced tables for all document types
    cursor.execute("""
    CREATE TABLE IF NOT EXISTS websites (
        id INTEGER PRIMARY KEY,
        sitemap_url TEXT UNIQUE,
        processed_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    """)
    cursor.execute("""
    CREATE TABLE IF NOT EXISTS pages (
        id INTEGER PRIMARY KEY,
        website_id INTEGER,
        url TEXT UNIQUE,
        title TEXT,
        content TEXT,
        processed_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (website_id) REFERENCES websites(id)
    )
    """)
    # New table for local documents
    cursor.execute("""
    CREATE TABLE IF NOT EXISTS documents (
        id INTEGER PRIMARY KEY,
        file_path TEXT UNIQUE,
        file_name TEXT,
        file_type TEXT,
        source_type TEXT,
        content TEXT,
        processed_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    """)
    # Create table for chat history
    cursor.execute("""
    CREATE TABLE IF NOT EXISTS chat_history (
        id INTEGER PRIMARY KEY,
        role TEXT,
        content TEXT,
        timestamp TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    """)
    conn.commit()
    return True


def save_vectorstore(vectorstore, save_path=DATA_DIR):
    """Save FAISS vectorstore to disk"""
    os.makedirs(save_path, exist_ok=True)
    vectorstore.save_local(f"{save_path}/faiss_index")


def load_vectorstore(load_path=DATA_DIR):
    """Load FAISS vectorstore from disk"""
    if os.path.exists(f"{load_path}/faiss_index"):
        embeddings = HuggingFaceEmbeddings()
        vectorstore = FAISS.load_local(
            f"{load_path}/faiss_index",
            embeddings,
            allow_dangerous_deserialization=True,
        )
        return vectorstore
    return None


def save_processed_content(sitemap_url, processed_urls, documents):
    """Save processed content to SQLite database"""
    conn = get_db_connection()
    cursor = conn.cursor()
    try:
        # Insert or update website record
        cursor.execute(
            "INSERT OR REPLACE INTO websites (sitemap_url, processed_date) VALUES (?, ?)",
            (sitemap_url, datetime.datetime.now()),
        )
        website_id = (
            cursor.lastrowid
            or cursor.execute(
                "SELECT id FROM websites WHERE sitemap_url = ?", (sitemap_url,)
            ).fetchone()[0]
        )
        # Insert document records
        for doc, url in zip(documents, processed_urls):
            # Extract title if available in metadata
            title = doc.metadata.get("title", os.path.basename(url))
            cursor.execute(
                "INSERT OR REPLACE INTO pages (website_id, url, title, content, processed_date) VALUES (?, ?, ?, ?, ?)",
                (website_id, url, title, doc.page_content, datetime.datetime.now()),
            )
        conn.commit()
    except Exception as e:
        conn.rollback()
        st.error(f"Database error: {str(e)}")


def save_local_documents(documents, file_paths, file_type):
    """Save locally processed documents to SQLite database"""
    conn = get_db_connection()
    cursor = conn.cursor()
    try:
        for doc, file_path in zip(documents, file_paths):
            file_name = os.path.basename(file_path)
            cursor.execute(
                "INSERT OR REPLACE INTO documents (file_path, file_name, file_type, source_type, content, processed_date) VALUES (?, ?, ?, ?, ?, ?)",
                (
                    file_path,
                    file_name,
                    file_type,
                    "local_directory",
                    doc.page_content,
                    datetime.datetime.now(),
                ),
            )
        conn.commit()
    except Exception as e:
        conn.rollback()
        st.error(f"Database error: {str(e)}")


def load_processed_content(sitemap_url=None):
    """Load processed content from SQLite database"""
    conn = get_db_connection()
    cursor = conn.cursor()
    documents = []
    processed_urls = []
    # Load web pages
    if sitemap_url:
        cursor.execute(
            """
            SELECT p.url, p.title, p.content
            FROM pages p
            JOIN websites w ON p.website_id = w.id
            WHERE w.sitemap_url = ?
            """,
            (sitemap_url,),
        )
    else:
        cursor.execute("SELECT url, title, content FROM pages")
    results = cursor.fetchall()
    for url, title, content in results:
        doc = Document(
            page_content=content,
            metadata={"source": url, "url": url, "title": title, "type": "web"},
        )
        documents.append(doc)
        processed_urls.append(url)
    # Load local documents
    cursor.execute("SELECT file_path, file_name, content, file_type FROM documents")
    results = cursor.fetchall()
    for file_path, file_name, content, file_type in results:
        doc = Document(
            page_content=content,
            metadata={
                "source": file_path,
                "file_name": file_name,
                "file_type": file_type,
                "type": "local",
            },
        )
        documents.append(doc)
        processed_urls.append(file_path)
    return documents, processed_urls


def save_chat_history(chat_history):
    """Save chat history to database"""
    conn = get_db_connection()
    cursor = conn.cursor()
    try:
        # Clear existing history
        cursor.execute("DELETE FROM chat_history")
        # Insert new history
        for idx, message in enumerate(chat_history):
            cursor.execute(
                "INSERT INTO chat_history (id, role, content) VALUES (?, ?, ?)",
                (idx, message["role"], message["content"]),
            )
        conn.commit()
    except Exception as e:
        conn.rollback()
        st.error(f"Error saving chat history: {str(e)}")


def load_chat_history():
    """Load chat history from database"""
    conn = get_db_connection()
    cursor = conn.cursor()
    try:
        cursor.execute("SELECT role, content FROM chat_history ORDER BY id")
        results = cursor.fetchall()
        chat_history = []
        for role, content in results:
            chat_history.append({"role": role, "content": content})
        return chat_history
    except:
        return []


def process_powerpoint_file(file_path: str) -> Document:
    """Process a PowerPoint file and extract text content"""
    try:
        # First try with Docling (it might support PPTX)
        try:
            converter = DocumentConverter()
            result = converter.convert(file_path)
            title = os.path.basename(file_path)
            # Try to get markdown format first
            try:
                full_text = result.document.export_to_markdown()
                if full_text and full_text.strip():
                    return Document(
                        page_content=full_text,
                        metadata={
                            "source": file_path,
                            "file_name": title,
                            "file_type": "pptx",
                            "format": "markdown",
                        },
                    )
            except AttributeError:
                pass
            # Try other content extraction methods
            if hasattr(result.document, "text"):
                text_content = result.document.text
            elif hasattr(result.document, "get_text"):
                text_content = result.document.get_text()
            else:
                text_content = str(result.document)
            if text_content and text_content.strip():
                return Document(
                    page_content=text_content,
                    metadata={
                        "source": file_path,
                        "file_name": title,
                        "file_type": "pptx",
                    },
                )
        except Exception:
            # If Docling fails, fall back to python-pptx
            pass
        # Fallback to manual extraction using python-pptx
        prs = Presentation(file_path)
        text_content = []
        slide_count = 0
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            slide_text.append(f"\n--- Slide {slide_num} ---\n")
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                # Handle tables in slides
                if shape.shape_type == 19:  # Table
                    try:
                        table_text = []
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                row_text.append(cell.text.strip())
                            table_text.append(" | ".join(row_text))
                        slide_text.append("\n".join(table_text))
                    except:
                        pass
            if len(slide_text) > 1:  # More than just the slide header
                text_content.extend(slide_text)
                slide_count += 1
        if text_content:
            full_text = "\n".join(text_content)
            title = os.path.basename(file_path)
            return Document(
                page_content=full_text,
                metadata={
                    "source": file_path,
                    "file_name": title,
                    "file_type": "pptx",
                    "slide_count": slide_count,
                    "extraction_method": "python-pptx",
                },
            )
        return None
    except Exception as e:
        st.warning(f"Error processing PowerPoint file {file_path}: {str(e)}")
        return None


def process_excel_file(file_path: str) -> Document:
    """Process an Excel file with better structure preservation"""
    try:
        text_content = []
        sheet_count = 0
        try:
            excel_file = pd.ExcelFile(file_path, engine="openpyxl")
        except:
            try:
                excel_file = pd.ExcelFile(file_path, engine="xlrd")
            except:
                excel_file = pd.ExcelFile(file_path)
        for sheet_name in excel_file.sheet_names:
            try:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                if df.empty:
                    continue
                sheet_text = [f"\n=== SHEET: {sheet_name} ===\n"]
                # Add structured data representation
                df_clean = df.fillna("")
                # For name-related data, create a more structured format
                if any(
                    col.lower() in ["name", "first_name", "last_name", "full_name"]
                    for col in df_clean.columns
                ):
                    sheet_text.append("STRUCTURED DATA - PEOPLE RECORDS:")
                    sheet_text.append("-" * 40)
                    for idx, row in df_clean.iterrows():
                        row_parts = []
                        for col, val in row.items():
                            if str(val).strip():
                                row_parts.append(f"{col}: {val}")
                        if row_parts:
                            sheet_text.append(
                                f"Record {idx + 1}: {' | '.join(row_parts)}"
                            )
                else:
                    # Standard table format for non-name data
                    if not df_clean.columns.empty:
                        headers = " | ".join(str(col) for col in df_clean.columns)
                        sheet_text.append(f"COLUMNS: {headers}")
                        sheet_text.append("-" * len(headers))
                    for idx, row in df_clean.iterrows():
                        row_text = " | ".join(
                            str(val) for val in row.values if str(val).strip()
                        )
                        if row_text.strip():
                            sheet_text.append(f"Row {idx + 1}: {row_text}")
                sheet_text.append(
                    f"\n[SHEET SUMMARY: {len(df)} total rows, {len(df.columns)} columns]"
                )
                sheet_text.append("=" * 50)
                text_content.extend(sheet_text)
                sheet_count += 1
            except Exception as e:
                st.warning(f"Error reading sheet '{sheet_name}': {str(e)}")
                continue
        if text_content:
            full_text = "\n".join(text_content)
            title = os.path.basename(file_path)
            return Document(
                page_content=full_text,
                metadata={
                    "source": file_path,
                    "file_name": title,
                    "file_type": "xlsx",
                    "sheet_count": sheet_count,
                    "extraction_method": "structured_pandas",
                    "total_content_length": len(full_text),
                },
            )
        return None
    except Exception as e:
        st.warning(f"Error processing Excel file {file_path}: {str(e)}")
        return None


def process_document_file(file_path: str, file_type: str) -> Document:
    """Process a single document file using appropriate method with OCR fallback for PDFs"""
    try:
        # Handle different file types with specialized functions
        if file_type.lower() == "pptx":
            return process_powerpoint_file(file_path)
        elif file_type.lower() in ["xlsx", "xls"]:
            return process_excel_file(file_path)
        elif file_type.lower() == "pdf":
            # For PDFs, try Docling first, then OCR if needed
            return process_pdf_with_intelligent_fallback(file_path)
        else:
            # Use original Docling approach for other formats
            converter = DocumentConverter()
            result = converter.convert(file_path)
            title = os.path.basename(file_path)

            # Try to get the full text as markdown
            try:
                full_text = result.document.export_to_markdown()
                if full_text and full_text.strip():
                    return Document(
                        page_content=full_text,
                        metadata={
                            "source": file_path,
                            "file_name": title,
                            "file_type": file_type,
                            "format": "markdown",
                        },
                    )
            except AttributeError:
                pass

            # If no markdown export, try to access content directly
            if hasattr(result.document, "text"):
                text_content = result.document.text
            elif hasattr(result.document, "get_text"):
                text_content = result.document.get_text()
            elif hasattr(result, "text"):
                text_content = result.text
            else:
                text_content = str(result.document)

            if text_content and text_content.strip():
                return Document(
                    page_content=text_content,
                    metadata={
                        "source": file_path,
                        "file_name": title,
                        "file_type": file_type,
                    },
                )
            return None
    except Exception as e:
        st.warning(f"Error processing file {file_path}: {str(e)}")
        return None


def process_pdf_with_intelligent_fallback(file_path: str) -> Document:
    """Process PDF with Docling first, then OCR fallback if needed"""
    try:
        st.write(f"**Processing PDF**: {os.path.basename(file_path)}")

        # First try regular Docling extraction
        docling_success = False
        docling_text = ""

        try:
            converter = DocumentConverter()
            result = converter.convert(file_path)

            # Try to get markdown format first
            try:
                full_text = result.document.export_to_markdown()
                if full_text and full_text.strip():
                    docling_text = full_text
            except AttributeError:
                pass

            # If no markdown, try other content extraction methods
            if not docling_text:
                if hasattr(result.document, "text"):
                    docling_text = result.document.text
                elif hasattr(result.document, "get_text"):
                    docling_text = result.document.get_text()
                else:
                    docling_text = str(result.document)

            # NOW DO PROPER QUALITY CHECK
            if docling_text and docling_text.strip():
                # Check multiple indicators of poor extraction:
                text_length = len(docling_text.strip())
                image_tag_count = docling_text.count("<!-- image -->")
                image_placeholder_count = docling_text.count("<image>")

                # Calculate ratio of actual text vs placeholders
                placeholder_chars = (image_tag_count * 15) + (
                    image_placeholder_count * 7
                )  # Approximate lengths
                actual_text_ratio = (
                    (text_length - placeholder_chars) / text_length
                    if text_length > 0
                    else 0
                )

                # Determine if extraction was successful based on multiple criteria:
                if (
                    text_length > 200  # Minimum meaningful length
                    and actual_text_ratio > 0.3  # At least 30% actual text
                    and image_tag_count
                    < (text_length / 100)  # Not too many image tags relative to content
                    and not (
                        docling_text.strip() == "<!-- image -->"
                        or docling_text.strip() == "<image>"
                    )
                ):  # Not just a single image tag
                    docling_success = True
                    st.write("✅ **Regular text extraction successful**")
                    st.write(
                        f"   📊 Extracted {text_length} characters with {actual_text_ratio:.1%} actual text"
                    )
                else:
                    st.write("⚠️ **Regular extraction produced poor quality results**")
                    st.write(
                        f"   📊 Length: {text_length}, Image tags: {image_tag_count}, Text ratio: {actual_text_ratio:.1%}"
                    )
            else:
                st.write("⚠️ **Regular extraction produced no text**")

        except Exception as e:
            st.write(f"⚠️ **Regular extraction failed with error**: {str(e)}")

        # If Docling was successful, return the result
        if docling_success and docling_text:
            return Document(
                page_content=docling_text,
                metadata={
                    "source": file_path,
                    "file_name": os.path.basename(file_path),
                    "file_type": "pdf",
                    "extraction_method": "docling",
                    "text_length": len(docling_text),
                },
            )

        # If regular extraction failed or gave poor results, try OCR
        st.write("🔍 **Attempting OCR extraction as fallback...**")
        ocr_result = process_pdf_with_ocr_internal(file_path)

        if ocr_result:
            st.write("✅ **OCR extraction successful**")
            return ocr_result
        else:
            st.write("❌ **Both regular and OCR extraction failed**")
            return None

    except Exception as e:
        st.warning(f"Error processing PDF {file_path}: {str(e)}")
        return None


def process_pdf_with_ocr_internal(file_path: str) -> Document:
    """Internal OCR processing function with better status reporting"""
    try:
        # Check if OCR libraries are available
        try:
            import pytesseract
            from pdf2image import convert_from_path
            from PIL import Image
        except ImportError as e:
            st.error(f"❌ **OCR libraries not installed**: {e}")
            st.write("Install with: `pip install pytesseract pillow pdf2image`")
            return None

        import tempfile
        import os

        # Convert PDF pages to images
        with tempfile.TemporaryDirectory() as temp_dir:
            try:
                # Convert PDF to images
                st.write("📄 **Converting PDF to images...**")
                pages = convert_from_path(file_path, dpi=300)
                st.write(f"📄 **Converted to {len(pages)} page images**")

                extracted_text = []
                successful_pages = 0

                # Create a progress indicator
                progress_placeholder = st.empty()

                for i, page in enumerate(pages):
                    progress_placeholder.write(
                        f"🔍 **Processing page {i + 1}/{len(pages)}...**"
                    )

                    try:
                        # Extract text using OCR
                        page_text = pytesseract.image_to_string(
                            page, config="--psm 6"
                        )  # Better OCR config
                        if page_text.strip():
                            extracted_text.append(f"\n--- Page {i + 1} ---\n")
                            extracted_text.append(page_text.strip())
                            successful_pages += 1
                    except Exception as e:
                        st.write(f"❌ **OCR failed for page {i + 1}**: {str(e)}")
                        continue

                progress_placeholder.empty()

                if extracted_text:
                    full_text = "\n".join(extracted_text)
                    st.write(f"✅ **OCR completed successfully!**")
                    st.write(f"   📊 Processed {successful_pages}/{len(pages)} pages")
                    st.write(f"   📊 Extracted {len(full_text)} characters")

                    return Document(
                        page_content=full_text,
                        metadata={
                            "source": file_path,
                            "file_name": os.path.basename(file_path),
                            "file_type": "pdf",
                            "extraction_method": "ocr",
                            "pages_processed": len(pages),
                            "successful_pages": successful_pages,
                            "text_length": len(full_text),
                        },
                    )
                else:
                    st.error("❌ **No text could be extracted via OCR from any page**")
                    return None

            except Exception as e:
                st.error(f"❌ **PDF to image conversion failed**: {str(e)}")
                return None

    except Exception as e:
        st.error(f"❌ **OCR processing failed**: {str(e)}")
        return None

    except ImportError:
        st.error(
            "❌ **OCR libraries not installed**. Please install: pip install pytesseract pillow pdf2image"
        )
        return None
    except Exception as e:
        st.error(f"❌ **OCR processing failed**: {str(e)}")
        return None


def process_local_directory(
    directory_path: str, file_extension: str, max_files: int = None
):
    """Process all files of a specific type in a directory"""
    if not os.path.exists(directory_path):
        st.error(f"Directory not found: {directory_path}")
        return [], []
    # Get all files with the specified extension
    pattern = os.path.join(directory_path, f"*.{file_extension}")
    file_paths = glob.glob(pattern)
    if not file_paths:
        st.warning(f"No .{file_extension} files found in {directory_path}")
        return [], []
    st.success(f"Found {len(file_paths)} .{file_extension} files")
    # Limit number of files if specified
    if max_files and max_files < len(file_paths):
        st.info(f"Processing only {max_files} files out of {len(file_paths)}")
        file_paths = file_paths[:max_files]
    all_documents = []
    processed_files = []
    # Process each file
    progress_bar = st.progress(0)
    status_text = st.empty()
    for idx, file_path in enumerate(file_paths):
        status_text.text(
            f"Processing file {idx + 1}/{len(file_paths)}: {os.path.basename(file_path)}"
        )
        try:
            document = process_document_file(file_path, file_extension)
            if document:
                all_documents.append(document)
                processed_files.append(file_path)
            progress_bar.progress((idx + 1) / len(file_paths))
            # Add small delay to prevent overwhelming the system
            time.sleep(0.1)
        except Exception as e:
            st.warning(f"Failed to process {file_path}: {str(e)}")
            continue
    progress_bar.empty()
    status_text.empty()
    if not all_documents:
        st.error(f"No documents could be extracted from .{file_extension} files")
    return all_documents, processed_files


def get_sitemap_urls(sitemap_url: str) -> List[str]:
    """
    Extracts all URLs from a sitemap XML file.
    Handles both regular sitemaps and sitemap indexes.
    """
    try:
        response = requests.get(sitemap_url, timeout=30)
        response.raise_for_status()
        root = ET.fromstring(response.content)
        # Define XML namespaces
        namespaces = {
            "sm": "http://www.sitemaps.org/schemas/sitemap/0.9",
            "xhtml": "http://www.w3.org/1999/xhtml",
        }
        urls = []
        # Check if this is a sitemap index (contains other sitemaps)
        sitemap_tags = root.findall(".//sm:sitemap", namespaces)
        if sitemap_tags:
            # This is a sitemap index
            for sitemap_tag in sitemap_tags:
                loc_tag = sitemap_tag.find("sm:loc", namespaces)
                if loc_tag is not None and loc_tag.text:
                    # Recursively get URLs from each child sitemap
                    child_sitemap_url = loc_tag.text.strip()
                    child_urls = get_sitemap_urls(child_sitemap_url)
                    urls.extend(child_urls)
        else:
            # This is a regular sitemap
            url_tags = root.findall(".//sm:url", namespaces)
            for url_tag in url_tags:
                loc_tag = url_tag.find("sm:loc", namespaces)
                if loc_tag is not None and loc_tag.text:
                    urls.append(loc_tag.text.strip())
        return urls
    except Exception as e:
        st.error(f"Error processing sitemap: {str(e)}")
        return []


def fetch_html_content(url: str) -> str:
    """Fetch HTML content from a URL"""
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
        }
        response = requests.get(url, headers=headers, timeout=30)
        response.raise_for_status()
        return response.text
    except Exception as e:
        st.warning(f"Error fetching HTML from {url}: {str(e)}")
        return ""


def process_url_with_docling(url: str, temp_dir: str) -> Document:
    """Process a URL's HTML content with Docling"""
    html_content = fetch_html_content(url)
    if not html_content:
        return None
    # Create a temporary HTML file
    url_filename = url.replace(":", "_").replace("/", "_").replace(".", "_") + ".html"
    temp_file_path = os.path.join(temp_dir, url_filename)
    with open(temp_file_path, "w", encoding="utf-8") as f:
        f.write(html_content)
    # Process with Docling
    try:
        converter = DocumentConverter()
        result = converter.convert(temp_file_path)
        # Extract page title if possible
        title = url
        try:
            title_match = re.search("<title>(.*?)</title>", html_content, re.IGNORECASE)
            if title_match:
                title = title_match.group(1)
        except:
            pass
            # Try to get the full text as markdown
        try:
            full_text = result.document.export_to_markdown()
            if full_text and full_text.strip():
                return Document(
                    page_content=full_text,
                    metadata={
                        "source": url,
                        "url": url,
                        "title": title,
                        "format": "markdown",
                    },
                )
        except AttributeError:
            pass
        # If no markdown export, try to access content directly
        if hasattr(result.document, "text"):
            text_content = result.document.text
        elif hasattr(result.document, "get_text"):
            text_content = result.document.get_text()
        elif hasattr(result, "text"):
            text_content = result.text
        else:
            text_content = str(result.document)
        if text_content and text_content.strip():
            return Document(
                page_content=text_content,
                metadata={
                    "source": url,
                    "url": url,
                    "title": title,
                },
            )
        # If still no content, try to access raw content
        if hasattr(result, "content"):
            content = result.content
            if content and str(content).strip():
                return Document(
                    page_content=str(content),
                    metadata={
                        "source": url,
                        "url": url,
                        "title": title,
                    },
                )
        return None
    except Exception as e:
        st.warning(f"Error processing URL with Docling: {url} - {str(e)}")
        return None
    finally:
        # Clean up the temp file
        if os.path.exists(temp_file_path):
            try:
                os.remove(temp_file_path)
            except:
                pass


def process_sitemap(sitemap_url: str, max_urls: int = None):
    """Process a sitemap and extract content from all URLs"""
    with st.spinner("Fetching sitemap URLs..."):
        urls = get_sitemap_urls(sitemap_url)
        if not urls:
            st.error("No URLs found in the sitemap.")
            return [], []
        st.success(f"Found {len(urls)} URLs in the sitemap.")
        if max_urls and max_urls < len(urls):
            st.info(f"Processing only {max_urls} URLs out of {len(urls)}.")
            urls = urls[:max_urls]
    # Create a temporary directory for HTML files
    temp_dir = os.path.join(working_dir, "temp_html")
    os.makedirs(temp_dir, exist_ok=True)
    all_documents = []
    processed_urls = []
    # Process each URL
    progress_bar = st.progress(0)
    status_text = st.empty()
    for idx, url in enumerate(urls):
        status_text.text(f"Processing URL {idx + 1}/{len(urls)}: {url}")
        try:
            # Add some delay to avoid overloading the server
            time.sleep(random.uniform(0.5, 2.0))
            document = process_url_with_docling(url, temp_dir)
            if document:
                all_documents.append(document)
                processed_urls.append(url)
            progress_bar.progress((idx + 1) / len(urls))
        except Exception as e:
            st.warning(f"Failed to process {url}: {str(e)}")
            continue
    progress_bar.empty()
    status_text.empty()
    # Clean up the temp directory
    try:
        import shutil

        shutil.rmtree(temp_dir)
    except:
        pass
    if not all_documents:
        st.error("No documents could be extracted from the sitemap URLs.")
    return all_documents, processed_urls


def setup_vectorstore(documents):
    embeddings = HuggingFaceEmbeddings()
    # Different splitting strategy for different document types
    doc_chunks = []
    for doc in documents:
        if doc.metadata.get("file_type") in ["xlsx", "xls"]:
            # For Excel files, use larger chunks to preserve record integrity
            text_splitter = CharacterTextSplitter(
                separator="\n",
                chunk_size=2000,  # Larger chunks for tabular data
                chunk_overlap=100,  # Smaller overlap
            )
        else:
            # Standard splitting for other documents
            text_splitter = CharacterTextSplitter(
                separator="\n",
                chunk_size=1000,
                chunk_overlap=200,
            )
        chunks = text_splitter.split_documents([doc])
        doc_chunks.extend(chunks)
    vectorstore = FAISS.from_documents(doc_chunks, embeddings)
    return vectorstore


def create_chain(vectorstore):
    """Original create_chain function"""
    llm = ChatGroq(
        model="llama-3.3-70b-versatile",
        temperature=0,
        max_tokens=4096,
    )
    retriever = vectorstore.as_retriever()
    memory = ConversationBufferMemory(
        llm=llm, output_key="answer", memory_key="chat_history", return_messages=True
    )
    chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=retriever,
        chain_type="map_reduce",
        memory=memory,
        verbose=True,
    )
    return chain


def create_chain_with_custom_retriever(vectorstore, k=20):
    """Create a chain with a custom retriever that returns more results"""
    llm = ChatGroq(
        model="llama-3.3-70b-versatile",
        temperature=0,
        max_tokens=4096,
    )
    # Create a custom retriever that returns more documents
    retriever = vectorstore.as_retriever(
        search_type="similarity",
        search_kwargs={"k": k},  # Return top k most similar chunks
    )
    memory = ConversationBufferMemory(
        llm=llm, output_key="answer", memory_key="chat_history", return_messages=True
    )
    chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=retriever,
        chain_type="map_reduce",  # Better for handling multiple documents
        memory=memory,
        verbose=True,
        return_source_documents=True,  # This helps with debugging
    )
    return chain


def initialize_storage():
    """Initialize storage and load existing data"""
    # Set up SQLite database
    setup_database()
    # Try to load existing vectorstore
    vectorstore = load_vectorstore(DATA_DIR)
    if vectorstore:
        st.session_state.vectorstore = vectorstore
        # Load URLs from database to session state
        documents, urls = load_processed_content()
        st.session_state.processed_urls = urls
        # Load chat history if available
        chat_history = load_chat_history()
        if chat_history:
            st.session_state.chat_history = chat_history
        # st.sidebar.success(f"Loaded {len(urls)} previously processed items")
    return True


def get_content_statistics():
    """Get statistics about processed content"""
    conn = get_db_connection()
    cursor = conn.cursor()
    stats = {}
    # Count web pages
    cursor.execute("SELECT COUNT(*) FROM pages")
    stats["web_pages"] = cursor.fetchone()[0]
    # Count documents by type
    cursor.execute("SELECT file_type, COUNT(*) FROM documents GROUP BY file_type")
    doc_counts = cursor.fetchall()
    for file_type, count in doc_counts:
        stats[f"{file_type}_files"] = count
    return stats


def on_shutdown():
    close_db_connection()


# Register shutdown handler
import atexit

atexit.register(on_shutdown)

##################### QUERY CONFIGURATION (begin) #################


def handle_improved_query_routing(query: str, vectorstore) -> str:
    """Improved query routing with clearer logic and exceptions"""
    try:
        query_lower = query.lower()
        # 1. EXACT name listing queries (very specific)
        if any(
            phrase in query_lower
            for phrase in [
                "list all names",
                "show all names",
                "all names in database",
                "list the names",
                "what names are in",
                "list all people named",
                "list people named",
                "show people named",
                "all people named",
                "people with the first name",
                "people with the last name",
                "list all the people named",
                "show all people named",
            ]
        ):
            # st.write("**Debug: Routing to name listing handler**")
            if any(
                phrase in query_lower
                for phrase in [
                    "list all names",
                    "show all names",
                    "all names in database",
                    "list the names",
                    "what names are in",
                ]
            ):
                names = get_all_names_from_vectorstore_improved(vectorstore)
                if names:
                    return (
                        f"Here are all {len(names)} names in the database:\n\n"
                        + "\n".join([f"{i}. {name}" for i, name in enumerate(names, 1)])
                    )
                else:
                    return "I couldn't find any names in the database."
            else:
                return handle_name_listing_query(query, vectorstore)
        # 2. Name count queries
        elif any(
            phrase in query_lower
            for phrase in ["how many names", "count names", "number of names"]
        ):
            names = get_all_names_from_vectorstore_improved(vectorstore)
            return f"There are {len(names)} names in the database."
        # 3. NEW: "Most recent" queries - HANDLE BEFORE general list queries
        elif "most recent" in query_lower:
            # st.write("**Debug: Routing to most recent handler**")
            # Get Excel content
            all_content = []
            if hasattr(vectorstore, "docstore") and hasattr(
                vectorstore, "index_to_docstore_id"
            ):
                for i in range(len(vectorstore.index_to_docstore_id)):
                    doc_id = vectorstore.index_to_docstore_id.get(i)
                    if doc_id and doc_id in vectorstore.docstore._dict:
                        doc = vectorstore.docstore._dict[doc_id]
                        if doc.metadata.get("file_type") in ["xlsx", "xls"]:
                            all_content.append(doc.page_content)
            if all_content:
                combined_content = "\n".join(all_content)
                return find_most_recent_client(query, combined_content)
            else:
                return "I couldn't find any Excel data to search through."
        # 4. NEW: Queries that start with "list" or "list all" - use Excel data (AFTER most recent check)
        elif query_lower.startswith(("list ", "list all ")):
            # st.write("**Debug: List query detected - routing to Excel search**")
            return handle_list_query(query, vectorstore)
        # 5. Document queries (agreements, summaries, etc.) - NON-EXCEL only
        elif any(
            phrase in query_lower
            for phrase in [
                "summarize",
                "summary of",
                "agreement",
                "contract",
                "document",
            ]
        ):
            # st.write("**Debug: Routing to document handler**")
            return handle_document_query(query, vectorstore)
        # 6. NEW: Exception queries - "I need all information on" or "Tell me about"
        elif any(
            phrase in query_lower
            for phrase in [
                "i need all information on",
                "tell me about",
                "give me all information on",
                "show me all information on",
            ]
        ):
            # st.write("**Debug: Full information request detected**")
            return handle_full_information_query(query, vectorstore)
        # 7. Person-specific queries - CHECK IF NAME EXISTS IN DATABASE FIRST
        else:
            # Extract potential names from the query
            words = query.split()
            skip_words = {
                "what",
                "is",
                "the",
                "phone",
                "number",
                "email",
                "address",
                "contact",
                "info",
                "information",
                "about",
                "tell",
                "me",
                "give",
                "show",
                "find",
                "get",
                "his",
                "her",
                "their",
                "of",
                "for",
                "on",
                "in",
                "at",
                "to",
                "from",
                "with",
                "and",
                "how",
                "can",
                "you",
                "please",
                "people",
                "named",
                "do",
                "know",
                "who",
                "date",
                "record",
                "case",
                "type",
            }
            potential_names = []
            i = 0
            while i < len(words):
                word = (
                    words[i]
                    .replace(",", "")
                    .replace(".", "")
                    .replace("?", "")
                    .replace("'s", "")
                )
                if (
                    word
                    and len(word) > 1
                    and word.isalpha()
                    and word.lower() not in skip_words
                ):
                    # Try to build a full name
                    name_parts = [word.capitalize()]
                    j = i + 1
                    while j < len(words) and j < i + 3:  # Max 3 parts
                        next_word = (
                            words[j]
                            .replace(",", "")
                            .replace(".", "")
                            .replace("?", "")
                            .replace("'s", "")
                        )
                        if (
                            next_word
                            and len(next_word) > 1
                            and next_word.isalpha()
                            and next_word.lower() not in skip_words
                        ):
                            name_parts.append(next_word.capitalize())
                            j += 1
                        else:
                            break
                    if len(name_parts) >= 2:
                        potential_names.append(" ".join(name_parts))
                        i = j
                    else:
                        potential_names.append(name_parts[0])
                        i += 1
                else:
                    i += 1
            # Check if any potential names exist in the database
            if potential_names:
                all_names = get_all_names_from_vectorstore_improved(vectorstore)
                name_found_in_db = False
                for search_name in potential_names:
                    for db_name in all_names:
                        # Check for exact match or if both first and last name appear
                        if search_name.lower() == db_name.lower():
                            name_found_in_db = True
                            break
                        # Check if it's a full name (first + last) that exists in DB
                        elif len(search_name.split()) >= 2:
                            search_parts = search_name.lower().split()
                            db_parts = db_name.lower().split()
                            if (
                                len(db_parts) >= 2
                                and search_parts[0] in db_parts
                                and search_parts[-1] in db_parts
                            ):
                                name_found_in_db = True
                                break
                    if name_found_in_db:
                        break
                if name_found_in_db:
                    # Name exists in database - use person-specific handler (Excel only)
                    # st.write("**Debug: Name found in database - routing to person handler**")
                    return handle_specific_person_query(query, vectorstore)
                else:
                    # Name not in database - use regular conversation chain (non-Excel documents)
                    # st.write("**Debug: Name not found in database - using regular chain**")
                    if "conversation_chain" in st.session_state:
                        response = st.session_state.conversation_chain(
                            {"question": query}
                        )
                        return response["answer"]
                    else:
                        return "I don't have enough context to answer that question."
            else:
                # No names detected - use regular conversation chain
                # st.write("**Debug: No names detected - using regular chain**")
                if "conversation_chain" in st.session_state:
                    response = st.session_state.conversation_chain({"question": query})
                    return response["answer"]
                else:
                    return "I don't have enough context to answer that question."
    except Exception as e:
        return f"I encountered an error while processing your question: {str(e)}"


def handle_name_listing_query(query: str, vectorstore) -> str:
    """Handle queries asking for lists of people with specific names"""
    try:
        # Extract the target name from the query
        query_lower = query.lower()
        words = query.split()
        # Find the name they're asking about
        target_name = None
        name_triggers = ["named", "name", "first", "last"]
        for i, word in enumerate(words):
            if word.lower() in name_triggers and i + 1 < len(words):
                # Get the next word as the target name
                target_name = (
                    words[i + 1].replace(",", "").replace(".", "").capitalize()
                )
                break
        if not target_name:
            return "I couldn't identify which name you're looking for. Please try 'list people named Robert' or 'people with first name John'."
        # st.write(f"**Debug: Looking for people named '{target_name}'**")
        # Get all names from database
        all_names = get_all_names_from_vectorstore_improved(vectorstore)
        # Find matches
        matching_names = []
        for db_name in all_names:
            name_parts = db_name.split()
            # Check first name match
            if name_parts[0].lower() == target_name.lower():
                matching_names.append(db_name)
            # Check last name match
            elif len(name_parts) > 1 and any(
                part.lower() == target_name.lower() for part in name_parts[1:]
            ):
                matching_names.append(db_name)
        if matching_names:
            result = (
                f"Found {len(matching_names)} people with the name '{target_name}':\n\n"
            )
            for i, name in enumerate(matching_names, 1):
                result += f"{i}. {name}\n"
            return result
        else:
            # Look for similar names
            similar_names = []
            for db_name in all_names:
                if target_name.lower() in db_name.lower():
                    similar_names.append(db_name)
            if similar_names:
                result = f"No exact matches for '{target_name}', but found similar names:\n\n"
                for i, name in enumerate(similar_names[:20], 1):  # Show up to 20
                    result += f"{i}. {name}\n"
                return result
            else:
                return f"No people found with the name '{target_name}'."
    except Exception as e:
        return f"Error processing name listing query: {str(e)}"


def handle_specific_person_query(query: str, vectorstore) -> str:
    """Handle queries about specific people - ONLY for allowed question types (unless it's a full info request)"""
    try:
        import re

        query_lower = query.lower()
        # Check if this is a full information request (exception to restrictions)
        is_full_info_request = any(
            phrase in query_lower
            for phrase in [
                "i need all information on",
                "tell me about",
                "give me all information on",
                "show me all information on",
            ]
        )
        # Check if the question is asking for allowed information (only if not full info request)
        if not is_full_info_request:
            allowed_questions = [
                "phone",
                "email",
                "date",
                "record",
                "case type",
                "contact",
            ]
            if not any(allowed in query_lower for allowed in allowed_questions):
                return "I can only answer questions about a person's phone number, email address, date of record, or case type."
        # Rest of your existing handle_specific_person_query code...
        # [Keep all the existing code from the original function]
        # STEP 1: Extract potential names (case insensitive)
        words = query.split()
        potential_names = []
        # Skip obvious non-name words
        skip_words = {
            "what",
            "is",
            "the",
            "phone",
            "number",
            "email",
            "address",
            "contact",
            "info",
            "information",
            "about",
            "tell",
            "me",
            "give",
            "show",
            "find",
            "get",
            "his",
            "her",
            "their",
            "of",
            "for",
            "on",
            "in",
            "at",
            "to",
            "from",
            "with",
            "and",
            "how",
            "can",
            "you",
            "please",
        }
        # FIXED: Look for potential names regardless of case
        i = 0
        while i < len(words):
            word = (
                words[i]
                .replace(",", "")
                .replace(".", "")
                .replace("?", "")
                .replace("'s", "")
            )
            # Check if this could be a name (alphabetic, not a skip word, reasonable length)
            if (
                word
                and len(word) > 1
                and word.isalpha()
                and word.lower() not in skip_words
            ):
                # Try to build a full name by looking ahead
                name_parts = [word.capitalize()]  # Capitalize for consistency
                j = i + 1
                # Keep adding consecutive potential name words
                while j < len(words) and j < i + 3:  # Limit to 3 parts max
                    next_word = (
                        words[j]
                        .replace(",", "")
                        .replace(".", "")
                        .replace("?", "")
                        .replace("'s", "")
                    )
                    if (
                        next_word
                        and len(next_word) > 1
                        and next_word.isalpha()
                        and next_word.lower() not in skip_words
                    ):
                        name_parts.append(next_word.capitalize())
                        j += 1
                    else:
                        break
                # Add the name(s) we found
                if len(name_parts) >= 2:
                    # Add full name
                    potential_names.append(" ".join(name_parts))
                    # Also add individual parts for partial matching
                    for part in name_parts:
                        potential_names.append(part)
                    i = j  # Skip ahead
                elif len(name_parts) == 1:
                    potential_names.append(name_parts[0])
                    i += 1
                else:
                    i += 1
            else:
                i += 1
        # Remove duplicates and prioritize longer names (more specific)
        potential_names = list(set(potential_names))
        potential_names = sorted(
            potential_names, key=lambda x: (len(x.split()), len(x)), reverse=True
        )
        # st.write(f"**Debug: Looking for these names: {potential_names}**")
        if not potential_names:
            return "I couldn't identify any names in your query. Please try asking like 'What is John Smith's phone number?'"
        # Get ALL the names from the database
        all_names = get_all_names_from_vectorstore_improved(vectorstore)
        # st.write(f"**Debug: Total names in database: {len(all_names)}**")
        # Find matches with priority for exact full name matches
        matched_names = []
        for search_name in potential_names:
            # st.write(f"**Debug: Searching for '{search_name}'**")
            # First, try exact matches (case insensitive)
            for db_name in all_names:
                if search_name.lower() == db_name.lower():
                    matched_names.append(db_name)
                    # st.write(f"**Debug: EXACT MATCH FOUND: {db_name}**")
                    # If we found an exact match for a full name, we can stop here
                    if len(search_name.split()) > 1:
                        break
            # If we found a full name exact match, don't look for partial matches
            if matched_names and len(search_name.split()) > 1:
                break
        # Remove duplicates
        matched_names = list(set(matched_names))
        if not matched_names:
            # Try partial matching as fallback
            for search_name in potential_names:
                search_parts = search_name.lower().split()
                for db_name in all_names:
                    db_parts = db_name.lower().split()
                    if all(part in db_parts for part in search_parts):
                        matched_names.append(db_name)
                        # st.write(f"**Debug: PARTIAL MATCH: {db_name}**")
        if not matched_names:
            similar_names = [
                name
                for name in all_names
                if any(part.lower() in name.lower() for part in potential_names)
            ]
            return f"I couldn't find matches for: {potential_names}. Similar names: {similar_names[:10]}"
        # If we have multiple matches, prioritize exact matches
        if len(matched_names) == 1:
            # st.write(f"**Debug: Found specific match: {matched_names[0]}**")
            pass
        else:
            # st.write(f"**Debug: Found {len(matched_names)} matches: {matched_names}**")
            pass
        # Get the full records - ONLY FROM EXCEL FILES
        all_content = []
        if hasattr(vectorstore, "docstore") and hasattr(
            vectorstore, "index_to_docstore_id"
        ):
            for i in range(len(vectorstore.index_to_docstore_id)):
                doc_id = vectorstore.index_to_docstore_id.get(i)
                if doc_id and doc_id in vectorstore.docstore._dict:
                    doc = vectorstore.docstore._dict[doc_id]
                    if doc.metadata.get("file_type") in ["xlsx", "xls"]:
                        all_content.append(doc.page_content)
        combined_content = "\n".join(all_content)
        # Find the full records for the matched names
        found_records = []
        for name in matched_names:
            name_parts = name.split()
            if len(name_parts) >= 2:
                first_name = name_parts[0]
                last_name = " ".join(name_parts[1:])
                lines = combined_content.split("\n")
                for i, line in enumerate(lines):
                    if (
                        first_name.lower() in line.lower()
                        and last_name.lower() in line.lower()
                    ):
                        record_lines = []
                        for j in range(i, min(i + 10, len(lines))):
                            record_lines.append(lines[j])
                            if j > i and (
                                "Record " in lines[j]
                                and any(
                                    name_part in lines[j]
                                    for name_part in [first_name, last_name]
                                )
                            ):
                                break
                        found_records.append("\n".join(record_lines))
                        break
        if found_records:
            llm = ChatGroq(
                model="llama-3.3-70b-versatile", temperature=0, max_tokens=4096
            )
            records_text = "\n\n".join(found_records)
            # Different prompts based on whether it's a full info request or restricted
            if is_full_info_request:
                prompt = f"""
                Question: {query}
                Here are the database records:
                {records_text}
                
                Please provide ALL available information about the person(s) from the database records.
                Include phone numbers, email addresses, dates, case types, and any other details available.
                """
            else:
                prompt = f"""
                Question: {query}
                Here is the database record for {matched_names[0] if len(matched_names) == 1 else "the requested people"}:
                {records_text}
                
                IMPORTANT: Only answer questions about:
                - Phone number
                - Email address  
                - Date of record
                - Case type
                
                Please provide a direct answer to the specific question asked.
                If the question is not about phone number, email, date of record, or case type, 
                respond with: "I can only provide information about phone numbers, email addresses, dates of records, and case types."
                """
            response = llm.invoke(prompt)
            return response.content
        else:
            return f"Found matching names {matched_names} but couldn't retrieve their full records."
    except Exception as e:
        # st.write(f"**Debug: Error: {str(e)}**")
        return f"Error: {str(e)}"


def handle_list_query(query: str, vectorstore) -> str:
    """Handle queries that start with 'list' or 'list all' - use Excel data"""
    try:
        query_lower = query.lower()

        # Get all Excel content
        all_content = []
        if hasattr(vectorstore, "docstore") and hasattr(
            vectorstore, "index_to_docstore_id"
        ):
            for i in range(len(vectorstore.index_to_docstore_id)):
                doc_id = vectorstore.index_to_docstore_id.get(i)
                if doc_id and doc_id in vectorstore.docstore._dict:
                    doc = vectorstore.docstore._dict[doc_id]
                    if doc.metadata.get("file_type") in ["xlsx", "xls"]:
                        all_content.append(doc.page_content)

        if not all_content:
            return "I couldn't find any Excel data to search through."

        combined_content = "\n".join(all_content)

        # Handle "most recent" queries FIRST (before other case type queries)
        if "most recent" in query_lower:
            return find_most_recent_client(query, combined_content)

        # Handle specific list queries with direct processing
        elif "email" in query_lower:
            return extract_all_emails(combined_content)
        elif "case type" in query_lower and "with a case type of" not in query_lower:
            return extract_all_case_types(combined_content)
        elif "with a case type of" in query_lower:
            return find_clients_by_case_type(query, combined_content)
        else:
            # For other list queries, use LLM with chunking
            return process_list_query_with_chunking(query, combined_content)

    except Exception as e:
        return f"Error processing list query: {str(e)}"


def extract_all_emails(content: str) -> str:
    """Extract all email addresses from Excel content"""
    import re

    try:
        # Find all email addresses using regex
        email_pattern = r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b"
        emails = re.findall(email_pattern, content)

        # Remove duplicates and sort
        unique_emails = sorted(list(set(emails)))

        if unique_emails:
            result = f"Found {len(unique_emails)} unique email addresses:\n\n"
            for i, email in enumerate(unique_emails, 1):
                result += f"{i}. {email}\n"
            return result
        else:
            return "No email addresses found in the database."

    except Exception as e:
        return f"Error extracting emails: {str(e)}"


def extract_all_case_types(content: str) -> str:
    """Extract all case types from Excel content by analyzing person records"""
    try:
        lines = content.split("\n")
        case_types = set()
        # Since individual person queries work, let's find all names first
        # and then extract case types from their records
        # Method 1: Look for any line that contains "Case Type:" or similar patterns
        for line in lines:
            line_clean = line.strip()
            if not line_clean:
                continue
            # Look for various case type patterns
            patterns_to_check = [
                "case type:",
                "case type =",
                "type:",
                "type =",
                "case:",
                "case =",
            ]
            line_lower = line_clean.lower()
            for pattern in patterns_to_check:
                if pattern in line_lower:
                    # Extract everything after the pattern
                    pattern_index = line_lower.find(pattern)
                    case_type_part = line_clean[pattern_index + len(pattern) :].strip()
                    # Clean up the case type
                    case_type_part = (
                        case_type_part.split(",")[0]
                        .split("|")[0]
                        .split("\t")[0]
                        .strip()
                    )
                    if case_type_part and len(case_type_part) > 1:
                        case_types.add(case_type_part)
        # Method 2: Use LLM to extract case types from sample records
        if len(case_types) < 10:  # If we didn't find many, try LLM approach
            # st.write("**Debug: Using LLM to extract case types from records**")
            # Take samples from different parts of the content
            content_length = len(content)
            sample_size = 3000
            samples = []
            # Take 5 samples from different positions
            for i in range(5):
                start_pos = (content_length // 5) * i
                end_pos = min(start_pos + sample_size, content_length)
                samples.append(content[start_pos:end_pos])
            llm = ChatGroq(
                model="llama-3.3-70b-versatile", temperature=0, max_tokens=1000
            )
            for i, sample in enumerate(samples):
                prompt = f"""
                Here is a sample of database records:
                {sample}
                
                Extract ONLY the case types mentioned in these records. 
                Look for fields like "Case Type:", "Type:", or similar.
                List each unique case type on a separate line.
                Do not include names, dates, or other information - only case types.
                If no case types found, respond with "NONE"
                """
                try:
                    response = llm.invoke(prompt)
                    if (
                        response.content.strip()
                        and "NONE" not in response.content.upper()
                    ):
                        # Parse the response and add to case_types
                        extracted_types = response.content.strip().split("\n")
                        for case_type in extracted_types:
                            case_type_clean = case_type.strip(" -•*123456789.")
                            if case_type_clean and len(case_type_clean) > 2:
                                case_types.add(case_type_clean)
                except Exception as e:
                    # st.write(f"**Debug: Error processing sample {i + 1}: {str(e)}**")
                    continue
        # Method 3: If still not enough, look for records and extract case types
        if len(case_types) < 5:
            # st.write("**Debug: Trying record-by-record extraction**")
            # Look for record patterns and extract case types
            record_blocks = []
            current_block = []
            for line in lines:
                if line.strip():
                    current_block.append(line)
                else:
                    if current_block:
                        record_blocks.append("\n".join(current_block))
                        current_block = []
            # Process blocks that look like they contain person data
            for block in record_blocks[:50]:  # Check first 50 blocks
                block_lower = block.lower()
                if any(
                    indicator in block_lower
                    for indicator in ["name", "phone", "email", "@"]
                ):
                    # This looks like a person record, extract case type
                    for line in block.split("\n"):
                        line_lower = line.lower()
                        if "case" in line_lower and "type" in line_lower:
                            # Try to extract the case type value
                            if ":" in line:
                                parts = line.split(":")
                                if len(parts) > 1:
                                    potential_case_type = parts[1].strip()
                                    if (
                                        potential_case_type
                                        and len(potential_case_type) > 2
                                    ):
                                        case_types.add(potential_case_type)
        # Clean up and deduplicate case types
        cleaned_case_types = set()
        for case_type in case_types:
            # Remove common false positives
            case_type_clean = case_type.strip()
            if (
                case_type_clean
                and len(case_type_clean) > 2
                and case_type_clean.lower()
                not in [
                    "case type",
                    "type",
                    "case",
                    "record",
                    "data",
                    "info",
                    "information",
                ]
            ):
                cleaned_case_types.add(case_type_clean)
        unique_case_types = sorted(list(cleaned_case_types))
        # st.write(f"**Debug: Found {len(unique_case_types)} case types: {unique_case_types}**")
        if unique_case_types:
            result = f"Found {len(unique_case_types)} case types:\n\n"
            for i, case_type in enumerate(unique_case_types, 1):
                result += f"{i}. {case_type}\n"
            return result
        else:
            return "No case types found in the database. The data structure might be different than expected."
    except Exception as e:
        return f"Error extracting case types: {str(e)}"


def find_clients_by_case_type(query: str, content: str) -> str:
    """Find clients with a specific case type"""
    try:
        # Check if user wants details (case types shown)
        if "show case type" in query.lower() or "with case type" in query.lower():
            return find_clients_by_case_type_with_details(query, content)
        else:
            return find_clients_by_case_type_comprehensive(query, content)
    except Exception as e:
        return f"Error finding clients by case type: {str(e)}"


def find_clients_by_case_type_comprehensive(query: str, content: str) -> str:
    """More precise case type matching and better name formatting"""
    try:
        # Extract the case type from the query
        query_lower = query.lower()
        case_type_start = query_lower.find("case type of ") + len("case type of ")
        case_type = query[case_type_start:].strip().lower()
        # st.write(f"**Debug: Looking for case type: '{case_type}'**")
        lines = content.split("\n")
        # st.write(f"**Debug: Processing {len(lines)} lines**")
        matching_clients = []
        # Process each line
        for i, line in enumerate(lines):
            line_clean = line.strip()
            # Skip empty lines
            if not line_clean:
                continue
            line_lower = line_clean.lower()
            # Must start with "Record" and contain case_type field
            if not line_clean.startswith("Record "):
                continue
            # More precise case type matching - look specifically in CASE_TYPE field
            case_type_match = False
            if "|" in line_clean:
                parts = line_clean.split("|")
                for part in parts:
                    part_clean = part.strip()
                    if part_clean.upper().startswith("CASE_TYPE:"):
                        case_type_value = part_clean.split(":", 1)[1].strip().lower()
                        if case_type in case_type_value:
                            case_type_match = True
                            break
            if case_type_match:
                # st.write(f"**Debug: Line {i + 1} contains case type '{case_type}'**")
                # Extract first and last name from this line
                first_name, last_name = extract_name_from_record_improved(line_clean)
                if first_name and last_name:
                    full_name = f"{first_name} {last_name}"
                    matching_clients.append(full_name)
                    # st.write(f"**Debug: Extracted name: '{full_name}'**")
                # else:
                # st.write(f"**Debug: Could not extract complete name from line {i + 1} (first: '{first_name}', last: '{last_name}')** ")
        # st.write(f"**Debug: Found {len(matching_clients)} matching clients**")
        # Remove duplicates and sort
        unique_clients = sorted(list(set(matching_clients)))
        if unique_clients:
            result = (
                f"Found {len(unique_clients)} clients with case type '{case_type}':\n\n"
            )
            for i, client in enumerate(unique_clients, 1):
                result += f"{i}. {client}\n"
            return result
        else:
            return f"No clients found with case type '{case_type}'"
    except Exception as e:
        return f"Error in comprehensive case type search: {str(e)}"


def extract_name_from_record_improved(record: str) -> tuple:
    """Extract first and last name with proper capitalization"""
    try:
        first_name = None
        last_name = None

        # The record is on a single line with pipe delimiters
        if "|" in record:
            parts = record.split("|")
            for part in parts:
                part = part.strip()

                # Handle first name (might be in format "Record X: FIRST_NAME: John")
                if "FIRST_NAME:" in part.upper():
                    # Split by FIRST_NAME: and take the part after it
                    first_name_split = part.upper().split("FIRST_NAME:")
                    if len(first_name_split) > 1:
                        first_name_raw = first_name_split[1].strip()
                        # Convert to proper case (capitalize first letter, rest lowercase)
                        first_name = first_name_raw.capitalize()

                # Handle last name (format "LAST_NAME: Smith")
                elif part.upper().startswith("LAST_NAME:"):
                    last_name_raw = part.split(":", 1)[1].strip()
                    # Convert to proper case
                    last_name = last_name_raw.capitalize()

        return first_name, last_name

    except Exception as e:
        return None, None


def find_clients_by_case_type_with_details(query: str, content: str) -> str:
    """Find clients and show their case types - corrected pipe-delimited format"""
    try:
        # Extract the case type from the query
        query_lower = query.lower()
        case_type_start = query_lower.find("case type of ") + len("case type of ")
        case_type = query[case_type_start:].strip().lower()

        st.write(f"**Debug: Looking for case type: '{case_type}' with details**")

        lines = content.split("\n")
        matching_clients = []

        # Process each line
        for i, line in enumerate(lines):
            line_clean = line.strip()

            if not line_clean:
                continue

            line_lower = line_clean.lower()

            # Check if this line contains our target case type
            if (
                case_type in line_lower
                and "case_type:" in line_lower
                and line_clean.startswith("Record ")
            ):
                # Extract details from this line
                first_name = None
                last_name = None
                actual_case_type = None

                if "|" in line_clean:
                    parts = line_clean.split("|")
                    for part in parts:
                        part = part.strip()

                        # Handle first name
                        if "FIRST_NAME:" in part.upper():
                            first_name_split = part.upper().split("FIRST_NAME:")
                            if len(first_name_split) > 1:
                                first_name = first_name_split[1].strip()

                        # Handle last name
                        elif part.upper().startswith("LAST_NAME:"):
                            last_name = part.split(":", 1)[1].strip()

                        # Handle case type
                        elif part.upper().startswith("CASE_TYPE:"):
                            actual_case_type = part.split(":", 1)[1].strip()

                if first_name and last_name and actual_case_type:
                    full_name = f"{first_name} {last_name}"
                    matching_clients.append(
                        {"name": full_name, "case_type": actual_case_type}
                    )

        # Remove duplicates based on name
        seen_names = set()
        unique_clients = []
        for client in matching_clients:
            if client["name"] not in seen_names:
                seen_names.add(client["name"])
                unique_clients.append(client)

        # Sort by name
        unique_clients.sort(key=lambda x: x["name"])

        if unique_clients:
            result = f"Found {len(unique_clients)} clients with case type containing '{case_type}':\n\n"
            for i, client in enumerate(unique_clients, 1):
                result += f"{i}. {client['name']} - {client['case_type']}\n"
            return result
        else:
            return f"No clients found with case type containing '{case_type}'"

    except Exception as e:
        return f"Error finding clients with details: {str(e)}"


def find_clients_by_case_type_llm_fallback_improved(
    case_type: str, content: str
) -> str:
    """Improved LLM fallback that processes smaller chunks more systematically"""
    try:
        # st.write(f"**Debug: Using improved LLM fallback for case type: '{case_type}'**")
        llm = ChatGroq(
            model="llama-3.3-70b-versatile",
            temperature=0,
            max_tokens=2000,  # Smaller response to focus on names only
        )
        # Use smaller chunks to ensure complete processing
        chunk_size = 8000  # Reduced chunk size
        content_chunks = [
            content[i : i + chunk_size] for i in range(0, len(content), chunk_size)
        ]
        all_clients = []
        for i, chunk in enumerate(content_chunks):
            # st.write(f"**Debug: Processing chunk {i + 1}/{len(content_chunks)} (length: {len(chunk)})**")
            prompt = f"""
            Find clients with case type containing "{case_type}".
            
            Data chunk {i + 1}/{len(content_chunks)}:
            {chunk}
            
            Rules:
            - Return ONLY client names (first and last name)
            - One name per line
            - No explanations or commentary
            - If no matches: return "NO_MATCHES"
            
            Names:
            """
            try:
                response = llm.invoke(prompt)
                response_lines = response.content.strip().split("\n")
                for line in response_lines:
                    line_clean = line.strip()
                    if "NO_MATCHES" in line_clean.upper():
                        continue
                    # Remove numbering
                    if ". " in line_clean:
                        line_clean = line_clean.split(". ", 1)[1]
                    # Validate it looks like a name
                    words = line_clean.split()
                    if (
                        len(words) >= 2
                        and len(words) <= 4
                        and all(
                            word.replace("'", "").replace("-", "").isalpha()
                            for word in words
                        )
                        and len(line_clean) < 50
                    ):
                        all_clients.append(line_clean)
                # st.write(f"**Debug: Chunk {i + 1} yielded {len([l for l in response_lines if l.strip() and 'NO_MATCHES' not in l.upper()])} names**")
            except Exception as e:
                # st.write(f"**Debug: Error in chunk {i + 1}: {str(e)}**")
                continue
        # Remove duplicates
        unique_clients = sorted(list(set(all_clients)))
        # st.write(f"**Debug: LLM fallback found {len(unique_clients)} unique clients**")
        if unique_clients:
            result = f"Found {len(unique_clients)} clients with case type containing '{case_type}':\n\n"
            for i, client in enumerate(unique_clients, 1):
                result += f"{i}. {client}\n"
            return result
        else:
            return f"No clients found with case type containing '{case_type}'"
    except Exception as e:
        return f"Error in LLM fallback: {str(e)}"


def extract_name_from_record(record: str) -> tuple:
    """Extract first and last name from a pipe-delimited record"""
    try:
        first_name = None
        last_name = None

        # The record is on a single line with pipe delimiters
        if "|" in record:
            parts = record.split("|")
            for part in parts:
                part = part.strip()

                # Handle first name (might be in format "Record X: FIRST_NAME: John")
                if "FIRST_NAME:" in part.upper():
                    # Split by FIRST_NAME: and take the part after it
                    first_name_split = part.upper().split("FIRST_NAME:")
                    if len(first_name_split) > 1:
                        first_name = first_name_split[1].strip()

                # Handle last name (format "LAST_NAME: Smith")
                elif part.upper().startswith("LAST_NAME:"):
                    last_name = part.split(":", 1)[1].strip()

        return first_name, last_name

    except Exception as e:
        return None, None


def process_case_type_in_chunks(case_type: str, content: str) -> str:
    """Process case type search in chunks when content is too large"""
    try:
        chunk_size = 15000
        chunks = [
            content[i : i + chunk_size] for i in range(0, len(content), chunk_size)
        ]

        llm = ChatGroq(model="llama-3.3-70b-versatile", temperature=0, max_tokens=4096)

        all_clients = []

        for i, chunk in enumerate(chunks):
            prompt = f"""
            Find clients with case type containing "{case_type}" in this data chunk {i + 1}/{len(chunks)}:
            
            {chunk}
            
            Return only the names, one per line. If no matches, return "NO_MATCHES".
            """

            response = llm.invoke(prompt)
            if "NO_MATCHES" not in response.content:
                chunk_clients = [
                    line.strip()
                    for line in response.content.split("\n")
                    if line.strip()
                ]
                all_clients.extend(chunk_clients)

        # Remove duplicates
        unique_clients = list(set(all_clients))

        if unique_clients:
            result = f"Found {len(unique_clients)} clients with case type containing '{case_type}':\n\n"
            for i, client in enumerate(sorted(unique_clients), 1):
                result += f"{i}. {client}\n"
            return result
        else:
            return f"No clients found with case type containing '{case_type}'"

    except Exception as e:
        return f"Error processing chunks: {str(e)}"


def find_most_recent_client(query: str, content: str) -> str:
    """Find the most recent client with specific criteria - updated for pipe format"""
    try:
        # Extract case type if specified
        case_type = None
        if "case type of" in query.lower():
            case_type_start = query.lower().find("case type of ") + len("case type of ")
            case_type = query[case_type_start:].strip().lower()
        # st.write(f"**Debug: Looking for most recent client with case type: '{case_type}'**")
        lines = content.split("\n")
        matching_records = []
        # Find all records with the specified case type
        for i, line in enumerate(lines):
            line_clean = line.strip()
            if not line_clean or not line_clean.startswith("Record "):
                continue
            line_lower = line_clean.lower()
            # More precise case type matching - look specifically in CASE_TYPE field
            case_type_match = False
            if "|" in line_clean:
                parts = line_clean.split("|")
                for part in parts:
                    part_clean = part.strip()
                    if part_clean.upper().startswith("CASE_TYPE:"):
                        case_type_value = part_clean.split(":", 1)[1].strip().lower()
                        if case_type and case_type in case_type_value:
                            case_type_match = True
                            break
            if case_type_match:
                # Extract name, date, and case type from this line
                first_name = None
                last_name = None
                date_found = None
                actual_case_type = None
                if "|" in line_clean:
                    parts = line_clean.split("|")
                    for part in parts:
                        part = part.strip()
                        # Extract first name
                        if "FIRST_NAME:" in part.upper():
                            first_name_split = part.upper().split("FIRST_NAME:")
                            if len(first_name_split) > 1:
                                first_name = first_name_split[1].strip().capitalize()
                        # Extract last name
                        elif part.upper().startswith("LAST_NAME:"):
                            last_name = part.split(":", 1)[1].strip().capitalize()
                        # Extract date
                        elif part.upper().startswith("DATE:"):
                            date_found = part.split(":", 1)[1].strip()
                        # Extract case type
                        elif part.upper().startswith("CASE_TYPE:"):
                            actual_case_type = part.split(":", 1)[1].strip()
                if first_name and last_name and date_found:
                    full_name = f"{first_name} {last_name}"
                    matching_records.append(
                        {
                            "name": full_name,
                            "date": date_found,
                            "case_type": actual_case_type,
                            "date_for_sorting": parse_date_for_sorting(date_found),
                        }
                    )
        # st.write(f"**Debug: Found {len(matching_records)} records with names and dates**")
        if matching_records:
            # Sort by date (most recent first)
            matching_records.sort(key=lambda x: x["date_for_sorting"], reverse=True)
            most_recent = matching_records[0]
            return f"Most recent client with case type '{case_type}':\n\n{most_recent['name']} (Date: {most_recent['date']}, Case Type: {most_recent['case_type']})"
        else:
            return f"No clients found with case type '{case_type}' that have complete name and date information."
    except Exception as e:
        return f"Error finding most recent client: {str(e)}"


def parse_date_for_sorting(date_str: str):
    """Convert date string to a format suitable for sorting"""
    try:
        import datetime

        # Try different date formats
        formats = [
            "%m/%d/%Y",
            "%m-%d-%Y",
            "%Y/%m/%d",
            "%Y-%m-%d",
            "%m/%d/%y",
            "%m-%d-%y",
        ]

        for fmt in formats:
            try:
                return datetime.datetime.strptime(date_str, fmt)
            except ValueError:
                continue

        # If no format works, return a default old date
        return datetime.datetime(1900, 1, 1)
    except:
        return datetime.datetime(1900, 1, 1)


def process_list_query_with_chunking(query: str, content: str) -> str:
    """Process list queries by chunking large content"""
    try:
        # Split content into manageable chunks
        chunk_size = 8000
        content_chunks = [
            content[i : i + chunk_size] for i in range(0, len(content), chunk_size)
        ]

        st.write(f"**Debug: Processing {len(content_chunks)} chunks**")

        llm = ChatGroq(model="llama-3.3-70b-versatile", temperature=0, max_tokens=4096)

        all_results = []

        for i, chunk in enumerate(content_chunks):
            st.write(f"**Debug: Processing chunk {i + 1}/{len(content_chunks)}**")

            prompt = f"""
            Question: {query}
            
            Here is a portion of the database content:
            {chunk}
            
            Instructions:
            - This is a list query, so provide a list format response
            - Search through this data portion to find records that match the criteria
            - Extract relevant information and format as a list
            - If no matches in this portion, respond with "No matches in this section"
            
            Answer:
            """

            response = llm.invoke(prompt)
            if (
                response.content.strip()
                and "No matches in this section" not in response.content
            ):
                all_results.append(response.content)

        if all_results:
            # Combine and deduplicate results
            combined_prompt = f"""
            Original query: {query}
            
            Here are results from different sections of the database:
            
            {chr(10).join([f"=== Section {i + 1} ==={chr(10)}{result}" for i, result in enumerate(all_results)])}
            
            Please combine these results, remove duplicates, and provide a final consolidated list.
            Count the total number of unique items found.
            """

            final_response = llm.invoke(combined_prompt)
            return final_response.content
        else:
            return f"No results found for: {query}"

    except Exception as e:
        return f"Error processing query with chunking: {str(e)}"


def handle_full_information_query(query: str, vectorstore) -> str:
    """Handle 'tell me about' or 'all information on' queries - combine Excel + other docs"""
    try:
        # Extract potential names from the query
        words = query.split()
        skip_words = {
            "i",
            "need",
            "all",
            "information",
            "on",
            "tell",
            "me",
            "about",
            "give",
            "show",
            "the",
            "what",
            "is",
            "and",
            "or",
            "a",
            "an",
            "to",
            "from",
            "with",
        }

        potential_names = []
        i = 0
        while i < len(words):
            word = (
                words[i]
                .replace(",", "")
                .replace(".", "")
                .replace("?", "")
                .replace("'s", "")
            )
            if (
                word
                and len(word) > 1
                and word.isalpha()
                and word.lower() not in skip_words
            ):
                name_parts = [word.capitalize()]
                j = i + 1
                while j < len(words) and j < i + 3:
                    next_word = (
                        words[j]
                        .replace(",", "")
                        .replace(".", "")
                        .replace("?", "")
                        .replace("'s", "")
                    )
                    if (
                        next_word
                        and len(next_word) > 1
                        and next_word.isalpha()
                        and next_word.lower() not in skip_words
                    ):
                        name_parts.append(next_word.capitalize())
                        j += 1
                    else:
                        break

                if len(name_parts) >= 2:
                    potential_names.append(" ".join(name_parts))
                    i = j
                else:
                    potential_names.append(name_parts[0])
                    i += 1
            else:
                i += 1

        if not potential_names:
            return "I couldn't identify any names in your query. Please specify who you want information about."

        st.write(f"Looking for full information on: {potential_names}")

        # Check if names exist in database
        all_names = get_all_names_from_vectorstore_improved(vectorstore)
        matched_names = []

        for search_name in potential_names:
            for db_name in all_names:
                if search_name.lower() == db_name.lower():
                    matched_names.append(db_name)
                    break
                elif len(search_name.split()) >= 2:
                    search_parts = search_name.lower().split()
                    db_parts = db_name.lower().split()
                    if (
                        len(db_parts) >= 2
                        and search_parts[0] in db_parts
                        and search_parts[-1] in db_parts
                    ):
                        matched_names.append(db_name)
                        break

        response_parts = []

        # If name found in database, get Excel data
        if matched_names:
            st.write("Name found - getting Excel data")
            excel_info = handle_specific_person_query_full(
                query, vectorstore, matched_names
            )
            response_parts.append("=== DATABASE INFORMATION ===")
            response_parts.append(excel_info)

        # Always get regular conversation chain data for full information requests
        st.write("Getting additional document information")
        if "conversation_chain" in st.session_state:
            regular_response = st.session_state.conversation_chain({"question": query})
            response_parts.append("=== ADDITIONAL DOCUMENT INFORMATION ===")
            response_parts.append(regular_response["answer"])

        if response_parts:
            return "\n\n".join(response_parts)
        else:
            return f"I couldn't find comprehensive information about {potential_names}."

    except Exception as e:
        return f"Error processing full information query: {str(e)}"


def handle_specific_person_query_full(
    query: str, vectorstore, matched_names: list
) -> str:
    """Get full person data without restrictions for 'tell me about' queries"""
    try:
        # Get all Excel content
        all_content = []
        if hasattr(vectorstore, "docstore") and hasattr(
            vectorstore, "index_to_docstore_id"
        ):
            for i in range(len(vectorstore.index_to_docstore_id)):
                doc_id = vectorstore.index_to_docstore_id.get(i)
                if doc_id and doc_id in vectorstore.docstore._dict:
                    doc = vectorstore.docstore._dict[doc_id]
                    if doc.metadata.get("file_type") in ["xlsx", "xls"]:
                        all_content.append(doc.page_content)

        combined_content = "\n".join(all_content)

        # Find records for matched names
        found_records = []
        for name in matched_names:
            name_parts = name.split()
            if len(name_parts) >= 2:
                first_name = name_parts[0]
                last_name = " ".join(name_parts[1:])
                lines = combined_content.split("\n")
                for i, line in enumerate(lines):
                    if (
                        first_name.lower() in line.lower()
                        and last_name.lower() in line.lower()
                    ):
                        record_lines = []
                        for j in range(i, min(i + 10, len(lines))):
                            record_lines.append(lines[j])
                            if j > i and (
                                "Record " in lines[j]
                                and any(
                                    name_part in lines[j]
                                    for name_part in [first_name, last_name]
                                )
                            ):
                                break
                        found_records.append("\n".join(record_lines))
                        break

        if found_records:
            llm = ChatGroq(
                model="llama-3.3-70b-versatile", temperature=0, max_tokens=4096
            )
            records_text = "\n\n".join(found_records)

            prompt = f"""
            Question: {query}
            Here are the database records:
            {records_text}
            
            Please provide ALL available information about the person(s) from the database records.
            Include phone numbers, email addresses, dates, case types, and any other details available.
            """

            response = llm.invoke(prompt)
            return response.content
        else:
            return f"Found matching names {matched_names} but couldn't retrieve their records from the database."

    except Exception as e:
        return f"Error retrieving full person information: {str(e)}"


def handle_document_query(query: str, vectorstore) -> str:
    """Handle document-specific queries like summarization with better content retrieval"""
    try:
        # Use a more comprehensive retrieval strategy for document queries
        retriever = vectorstore.as_retriever(
            search_type="similarity",
            search_kwargs={"k": 25},  # Get more chunks for comprehensive coverage
        )
        
        # Get relevant documents
        relevant_docs = retriever.get_relevant_documents(query)
        if not relevant_docs:
            return "I couldn't find any relevant documents for your query."
        
        # Filter out Excel files - only use PDF, DOCX, HTML, TXT files
        document_docs = []
        for doc in relevant_docs:
            file_type = doc.metadata.get("file_type", "").lower()
            if file_type not in ["xlsx", "xls"]:  # Exclude Excel files
                document_docs.append(doc)
        
        if not document_docs:
            return "I couldn't find any relevant non-Excel documents for your query."
        
        # Group documents by source and get more complete content
        docs_by_source = {}
        for doc in document_docs:
            source = doc.metadata.get("source", "unknown")
            file_name = doc.metadata.get("file_name", os.path.basename(source))
            if file_name not in docs_by_source:
                docs_by_source[file_name] = []
            docs_by_source[file_name].append(doc)
        
        # Combine content from relevant documents with larger chunks
        combined_content = []
        for file_name, docs in docs_by_source.items():
            combined_content.append(f"\n=== FROM DOCUMENT: {file_name} ===")
            for doc in docs:
                # Use full content instead of truncating to 1000 chars
                combined_content.append(doc.page_content)
        
        content_text = "\n".join(combined_content)
        
        # Process in chunks if content is too large
        max_content_length = 25000  # Adjust based on model limits
        
        if len(content_text) > max_content_length:
            return process_large_document_query(query, content_text, docs_by_source)
        else:
            return process_single_document_query(query, content_text)
            
    except Exception as e:
        return f"Error processing document query: {str(e)}"

def process_single_document_query(query: str, content_text: str) -> str:
    """Process document query when content fits in single request"""
    llm = ChatGroq(
        model="llama-3.3-70b-versatile",
        temperature=0,
        max_tokens=4096,
    )
    
    prompt = f"""
Based on the following document content, please answer this question: {query}

Document Content:
{content_text}

Instructions:
- Answer based ONLY on the provided document content
- If you're asked to summarize a specific agreement, provide all available details
- If you're asked for specific paragraphs (like ordered, adjudged, decreed), list them completely
- Be specific about which document you're referencing
- If you can't find the specific document mentioned, say so clearly
- Provide detailed information from the relevant document(s)
- Do NOT reference any Excel or database records

Answer:
"""
    
    response = llm.invoke(prompt)
    return response.content

def process_large_document_query(query: str, content_text: str, docs_by_source: dict) -> str:
    
    
    
    """Process document query when content is too large for single request"""
    llm = ChatGroq(
        model="llama-3.3-70b-versatile",
        temperature=0,
        max_tokens=4096,
    )
    
    # Process each document separately
    document_summaries = []
    
    for file_name, docs in docs_by_source.items():
        doc_content = "\n".join([doc.page_content for doc in docs])
        
        # Chunk large documents
        chunk_size = 15000
        if len(doc_content) > chunk_size:
            chunks = [doc_content[i:i + chunk_size] for i in range(0, len(doc_content), chunk_size)]
            chunk_results = []
            
            for i, chunk in enumerate(chunks):
                prompt = f"""
Based on this portion of document "{file_name}", answer: {query}

Document chunk {i+1}/{len(chunks)}:
{chunk}

Instructions:
- Focus on answering the specific question asked
- If this chunk contains relevant information, extract it completely
- If no relevant information in this chunk, respond with "No relevant information in this section"

Answer:
"""
                response = llm.invoke(prompt)
                if "No relevant information in this section" not in response.content:
                    chunk_results.append(response.content)
            
            if chunk_results:
                # Combine chunk results
                combined_prompt = f"""
These are results from different sections of document "{file_name}" for the query: {query}

{chr(10).join([f"Section {i+1}: {result}" for i, result in enumerate(chunk_results)])}

Please combine these results into a comprehensive answer, removing any duplicates.
"""
                final_response = llm.invoke(combined_prompt)
                document_summaries.append(f"From {file_name}:\n{final_response.content}")
        else:
            # Process entire document at once
            prompt = f"""
Based on document "{file_name}", answer: {query}

Document content:
{doc_content}

Instructions:
- Provide a complete and detailed answer
- Include all relevant information from this document

Answer:
"""
            response = llm.invoke(prompt)
            document_summaries.append(f"From {file_name}:\n{response.content}")
    
    # Combine all document results
    if len(document_summaries) > 1:
        final_prompt = f"""
Here are results from multiple documents for the query: {query}

{chr(10).join(document_summaries)}

Please provide a comprehensive final answer combining information from all documents.
"""
        final_response = llm.invoke(final_prompt)
        return final_response.content
    elif document_summaries:
        return document_summaries[0]
    else:
        return "No relevant information found in the documents."
###############################  Query Config End


def process_pdf_with_ocr(file_path: str) -> Document:
    """Process a PDF using OCR if regular text extraction fails"""
    try:
        import pytesseract
        from pdf2image import convert_from_path
        from PIL import Image
        import tempfile
        import os

        st.write(f"**Processing PDF with OCR**: {os.path.basename(file_path)}")

        # First try regular Docling extraction
        try:
            converter = DocumentConverter()
            result = converter.convert(file_path)

            if hasattr(result.document, "text"):
                text_content = result.document.text
            elif hasattr(result.document, "get_text"):
                text_content = result.document.get_text()
            else:
                text_content = str(result.document)

            # Check if we got meaningful text (more than just image tags)
            if len(text_content) > 500 and "<!-- image -->" not in text_content:
                st.write("✅ **Regular text extraction successful**")
                return Document(
                    page_content=text_content,
                    metadata={
                        "source": file_path,
                        "file_name": os.path.basename(file_path),
                        "file_type": "pdf",
                        "extraction_method": "docling",
                    },
                )
        except Exception as e:
            st.write(f"⚠️ **Regular extraction failed**: {str(e)}")

        st.write("🔍 **Attempting OCR extraction...**")

        # Convert PDF pages to images
        with tempfile.TemporaryDirectory() as temp_dir:
            try:
                # Convert PDF to images
                pages = convert_from_path(file_path, dpi=300)
                st.write(f"📄 **Converted to {len(pages)} images**")

                extracted_text = []

                for i, page in enumerate(pages):
                    st.write(f"🔍 **Processing page {i + 1}/{len(pages)}...**")

                    # Save page as temporary image
                    temp_image_path = os.path.join(temp_dir, f"page_{i + 1}.png")
                    page.save(temp_image_path, "PNG")

                    # Extract text using OCR
                    try:
                        page_text = pytesseract.image_to_string(page)
                        if page_text.strip():
                            extracted_text.append(f"\n--- Page {i + 1} ---\n")
                            extracted_text.append(page_text.strip())
                        else:
                            st.write(f"⚠️ **No text found on page {i + 1}**")
                    except Exception as e:
                        st.write(f"❌ **OCR failed for page {i + 1}**: {str(e)}")
                        continue

                if extracted_text:
                    full_text = "\n".join(extracted_text)
                    st.write(
                        f"✅ **OCR successful! Extracted {len(full_text)} characters**"
                    )

                    # Show preview
                    st.write("**OCR Text Preview (first 500 characters):**")
                    st.text_area("OCR Preview", full_text[:500], height=150)

                    return Document(
                        page_content=full_text,
                        metadata={
                            "source": file_path,
                            "file_name": os.path.basename(file_path),
                            "file_type": "pdf",
                            "extraction_method": "ocr",
                            "pages_processed": len(pages),
                            "text_length": len(full_text),
                        },
                    )
                else:
                    st.error("❌ **No text could be extracted via OCR**")
                    return None

            except Exception as e:
                st.error(f"❌ **PDF to image conversion failed**: {str(e)}")
                return None

    except ImportError:
        st.error(
            "❌ **OCR libraries not installed**. Please install: pip install pytesseract pillow pdf2image"
        )
        return None
    except Exception as e:
        st.error(f"❌ **OCR processing failed**: {str(e)}")
        return None


# Check authentication before showing the app
if not check_authentication():
    st.stop()

# Streamlit app configuration (only runs if authenticated)
st.set_page_config(page_title="Claudia-User", page_icon="🦙", layout="wide")

# Add logout button in the top right
col1, col2 = st.columns([4, 1])
with col1:
    st.title("CLAUDIA  🦙(LLAMA 3.3) - User")
with col2:
    # Show session info
    login_time = st.session_state.get("login_timestamp", 0)
    if login_time > 0:
        hours_elapsed = (time.time() - login_time) / 3600
        st.write(f"Welcome, {st.session_state.get('username', 'User')}")
        st.caption(f"Session: {24 - int(hours_elapsed)}h remaining")
    else:
        st.write(f"Welcome, {st.session_state.get('username', 'User')}")

    if st.button("Logout", type="secondary"):
        logout()

st.markdown(
    """
    <style>
        .reportview-container {
            margin-top: -2em;
        }
        #MainMenu {visibility: hidden;}
        .stDeployButton {display:none;}
        footer {visibility: hidden;}
        #stDecoration {display:none;}
    </style>
""",
    unsafe_allow_html=True,
)

# Initialize session state
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "processed_urls" not in st.session_state:
    st.session_state.processed_urls = []

# Initialize storage when app starts
if "storage_initialized" not in st.session_state:
    initialize_storage()
    st.session_state.storage_initialized = True

# Document loading options
# st.subheader("Load Content")

# # Sitemap input
# sitemap_col, sitemap_options_col = st.columns([3, 1])
# with sitemap_col:
#     sitemap_url = st.text_input(
#         "Enter website sitemap URL:",
#         placeholder="e.g., https://example.com/sitemap.xml",
#     )
# with sitemap_options_col:
#     max_urls = st.slider(
#         "Max URLs",
#         min_value=1,
#         max_value=100,
#         value=50,
#         help="Limit the number of URLs to process",
#     )

# # Local HTML directory input
# html_col, html_options_col = st.columns([3, 1])
# with html_col:
#     html_dir = st.text_input(
#         "Enter path to HTML files directory: /mnt/c/AI/Add/html",
#         placeholder="e.g., C:/Documents/html_files or /home/user/html_files",
#     )
# with html_options_col:
#     max_html_files = st.slider(
#         "Max HTML files",
#         min_value=1,
#         max_value=100,
#         value=50,
#         help="Limit the number of HTML files to process",
#     )

# # Local DOCX directory input
# docx_col, docx_options_col = st.columns([3, 1])
# with docx_col:
#     docx_dir = st.text_input(
#         "Enter path to DOCX files directory:  /mnt/c/AI/Add/docx",
#         placeholder="e.g., C:/Documents/docx_files or /home/user/docx_files",
#     )
# with docx_options_col:
#     max_docx_files = st.slider(
#         "Max DOCX files",
#         min_value=1,
#         max_value=100,
#         value=50,
#         help="Limit the number of DOCX files to process",
#     )

# # Local PDF directory input
# pdf_col, pdf_options_col = st.columns([3, 1])
# with pdf_col:
#     pdf_dir = st.text_input(
#         "Enter path to PDF files directory:  /mnt/c/AI/Add/pdf",
#         placeholder="e.g., C:/Documents/pdf_files or /home/user/pdf_files",
#     )
# with pdf_options_col:
#     max_pdf_files = st.slider(
#         "Max PDF files",
#         min_value=1,
#         max_value=100,
#         value=50,
#         help="Limit the number of PDF files to process",
#     )

# # Local PowerPoint directory input
# pptx_col, pptx_options_col = st.columns([3, 1])
# with pptx_col:
#     pptx_dir = st.text_input(
#         "Enter path to PowerPoint files directory: /mnt/c/AI/Add/pptx",
#         placeholder="e.g., C:/Documents/pptx_files or /home/user/pptx_files",
#     )
# with pptx_options_col:
#     max_pptx_files = st.slider(
#         "Max PPTX files",
#         min_value=1,
#         max_value=100,
#         value=50,
#         help="Limit the number of PowerPoint files to process",
#     )

# # Local Excel directory input
# xlsx_col, xlsx_options_col = st.columns([3, 1])
# with xlsx_col:
#     xlsx_dir = st.text_input(
#         "Enter path to Excel files directory: /mnt/c/AI/Add/xlsx",
#         placeholder="e.g., C:/Documents/xlsx_files or /home/user/xlsx_files",
#     )
# with xlsx_options_col:
#     max_xlsx_files = st.slider(
#         "Max Excel files",
#         min_value=1,
#         max_value=100,
#         value=50,
#         help="Limit the number of Excel files to process",
#     )

# # Process buttons
# col1, col2, col3, col4, col5, col6 = st.columns(6)
# with col1:
#     process_sitemap_button = st.button(
#         "Process Sitemap", type="primary", use_container_width=True
#     )
# with col2:
#     process_html_button = st.button(
#         "Process HTML", type="primary", use_container_width=True
#     )
# with col3:
#     process_docx_button = st.button(
#         "Process DOCX", type="primary", use_container_width=True
#     )
# with col4:
#     process_pdf_button = st.button(
#         "Process PDF", type="primary", use_container_width=True
#     )
# with col5:
#     process_pptx_button = st.button(
#         "Process PowerPoint", type="primary", use_container_width=True
#     )
# with col6:
#     process_xlsx_button = st.button(
#         "Process Excel", type="primary", use_container_width=True
#     )

# with col4:
#     # Add this test button for OCR Programs Installed
#     if st.button("🧪 Test Poppler", key="test_poppler"):
#         try:
#             import subprocess

#             # Test if poppler is accessible
#             result = subprocess.run(["pdftoppm", "-h"], capture_output=True, text=True)
#             if result.returncode == 0:
#                 st.success("✅ Poppler is installed and accessible!")
#             else:
#                 st.error("❌ Poppler not found in PATH")
#         except FileNotFoundError:
#             st.error("❌ Poppler not installed or not in PATH")
#             st.write("**Install instructions:**")
#             st.code("sudo apt install poppler-utils")
#         except Exception as e:
#             st.error(f"❌ Error testing Poppler: {e}")

# with col4:
#     if st.button("🧪 Test OCR", key="test_ocr_setup"):
#         try:
#             import pytesseract
#             import pdf2image
#             from PIL import Image

#             st.success("✅ All OCR libraries are installed!")
#             # Test Tesseract
#             version = pytesseract.get_tesseract_version()
#             st.write(f"📝 Tesseract version: {version}")
#         except ImportError as e:
#             st.error(f"❌ Missing library: {e}")
#             st.write("Install with: `pip install pytesseract pillow pdf2image`")
#         except Exception as e:
#             st.error(f"❌ Tesseract not found: {e}")
#             st.write("Install Tesseract: https://github.com/UB-Mannheim/tesseract/wiki")

# Process sitemap
# if process_sitemap_button and sitemap_url:
#     try:
#         documents, processed_urls = process_sitemap(sitemap_url, max_urls)
#         if documents:
#             st.success(
#                 f"Successfully processed {len(documents)} pages from {len(processed_urls)} URLs"
#             )
#             # Create or update vectorstore
#             if "vectorstore" not in st.session_state:
#                 st.session_state.vectorstore = setup_vectorstore(documents)
#             else:
#                 st.session_state.vectorstore.add_documents(documents)
#             # Save data to persistent storage
#             save_vectorstore(st.session_state.vectorstore)
#             save_processed_content(sitemap_url, processed_urls, documents)
#             st.session_state.processed_urls.extend(processed_urls)
#             # Show processed URLs
#             with st.expander("Processed URLs", expanded=True):
#                 for url in processed_urls:
#                     st.write(url)
#     except Exception as e:
#         st.error(f"Error processing sitemap: {str(e)}")

# # Process HTML files
# if process_html_button and html_dir:
#     try:
#         documents, processed_files = process_local_directory(
#             html_dir, "html", max_html_files
#         )
#         if documents:
#             st.success(f"Successfully processed {len(documents)} HTML files")
#             # Create or update vectorstore
#             if "vectorstore" not in st.session_state:
#                 st.session_state.vectorstore = setup_vectorstore(documents)
#             else:
#                 st.session_state.vectorstore.add_documents(documents)
#             # Save data to persistent storage
#             save_vectorstore(st.session_state.vectorstore)
#             save_local_documents(documents, processed_files, "html")
#             st.session_state.processed_urls.extend(processed_files)
#             # Show processed files
#             with st.expander("Processed HTML Files", expanded=True):
#                 for file_path in processed_files:
#                     st.write(os.path.basename(file_path))
#     except Exception as e:
#         st.error(f"Error processing HTML files: {str(e)}")

# # Process DOCX files
# if process_docx_button and docx_dir:
#     try:
#         documents, processed_files = process_local_directory(
#             docx_dir, "docx", max_docx_files
#         )
#         if documents:
#             st.success(f"Successfully processed {len(documents)} DOCX files")
#             # Create or update vectorstore
#             if "vectorstore" not in st.session_state:
#                 st.session_state.vectorstore = setup_vectorstore(documents)
#             else:
#                 st.session_state.vectorstore.add_documents(documents)
#             # Save data to persistent storage
#             save_vectorstore(st.session_state.vectorstore)
#             save_local_documents(documents, processed_files, "docx")
#             st.session_state.processed_urls.extend(processed_files)
#             # Show processed files
#             with st.expander("Processed DOCX Files", expanded=True):
#                 for file_path in processed_files:
#                     st.write(os.path.basename(file_path))
#     except Exception as e:
#         st.error(f"Error processing DOCX files: {str(e)}")

# # Process PDF files
# if process_pdf_button and pdf_dir:
#     try:
#         documents, processed_files = process_local_directory(
#             pdf_dir, "pdf", max_pdf_files
#         )
#         if documents:
#             st.success(f"Successfully processed {len(documents)} PDF files")
#             # Create or update vectorstore
#             if "vectorstore" not in st.session_state:
#                 st.session_state.vectorstore = setup_vectorstore(documents)
#             else:
#                 st.session_state.vectorstore.add_documents(documents)
#             # Save data to persistent storage
#             save_vectorstore(st.session_state.vectorstore)
#             save_local_documents(documents, processed_files, "pdf")
#             st.session_state.processed_urls.extend(processed_files)
#             # Show processed files
#             with st.expander("Processed PDF Files", expanded=True):
#                 for file_path in processed_files:
#                     st.write(os.path.basename(file_path))
#     except Exception as e:
#         st.error(f"Error processing PDF files: {str(e)}")

# # Process PowerPoint files
# if process_pptx_button and pptx_dir:
#     try:
#         documents, processed_files = process_local_directory(
#             pptx_dir, "pptx", max_pptx_files
#         )
#         if documents:
#             st.success(f"Successfully processed {len(documents)} PowerPoint files")
#             # Create or update vectorstore
#             if "vectorstore" not in st.session_state:
#                 st.session_state.vectorstore = setup_vectorstore(documents)
#             else:
#                 st.session_state.vectorstore.add_documents(documents)
#             # Save data to persistent storage
#             save_vectorstore(st.session_state.vectorstore)
#             save_local_documents(documents, processed_files, "pptx")
#             st.session_state.processed_urls.extend(processed_files)
#             # Show processed files
#             with st.expander("Processed PowerPoint Files", expanded=True):
#                 for file_path in processed_files:
#                     st.write(os.path.basename(file_path))
#     except Exception as e:
#         st.error(f"Error processing PowerPoint files: {str(e)}")

# # Process Excel files
# if process_xlsx_button and xlsx_dir:
#     try:
#         documents, processed_files = process_local_directory(
#             xlsx_dir, "xlsx", max_xlsx_files
#         )
#         if documents:
#             st.success(f"Successfully processed {len(documents)} Excel files")
#             # Create or update vectorstore
#             if "vectorstore" not in st.session_state:
#                 st.session_state.vectorstore = setup_vectorstore(documents)
#             else:
#                 st.session_state.vectorstore.add_documents(documents)
#             # Save data to persistent storage
#             save_vectorstore(st.session_state.vectorstore)
#             save_local_documents(documents, processed_files, "xlsx")
#             st.session_state.processed_urls.extend(processed_files)
#             # Show processed files
#             with st.expander("Processed Excel Files", expanded=True):
#                 for file_path in processed_files:
#                     st.write(os.path.basename(file_path))
#     except Exception as e:
#         st.error(f"Error processing Excel files: {str(e)}")


# # Add this as a completely separate section - put it after your other diagnostics
# st.write("---")
# st.subheader("🔍 OCR PDF Processing")

# # Simple file path input
# ocr_pdf_path = st.text_input(
#     "PDF file path for OCR:",
#     value="/mnt/c/AI/CRMRecords/agreement-S.pdf",
#     key="simple_ocr_path",
# )

# # Simple process button
# if st.button("Start OCR Processing", key="simple_ocr_start"):
#     if ocr_pdf_path and os.path.exists(ocr_pdf_path):
#         st.write(f"🔍 **Starting OCR on**: {os.path.basename(ocr_pdf_path)}")

#         try:
#             # Show that we're starting
#             with st.status("Processing PDF with OCR...", expanded=True) as status:
#                 st.write("📄 Converting PDF to images...")

#                 # Import here to check for errors
#                 try:
#                     import pytesseract
#                     import pdf2image
#                     from PIL import Image

#                     st.write("✅ OCR libraries loaded")
#                 except ImportError as e:
#                     st.error(f"❌ Missing libraries: {e}")
#                     st.stop()

#                 # Simple OCR processing
#                 st.write("🔄 Starting OCR extraction...")

#                 # Convert PDF to images
#                 pages = pdf2image.convert_from_path(ocr_pdf_path, dpi=200)
#                 st.write(f"📄 Converted to {len(pages)} page images")

#                 # Extract text from each page
#                 all_text = []
#                 for i, page in enumerate(pages):
#                     st.write(f"🔍 Processing page {i + 1}/{len(pages)}")
#                     try:
#                         page_text = pytesseract.image_to_string(page)
#                         if page_text.strip():
#                             all_text.append(
#                                 f"\n--- Page {i + 1} ---\n{page_text.strip()}"
#                             )
#                         else:
#                             st.write(f"⚠️ No text on page {i + 1}")
#                     except Exception as e:
#                         st.write(f"❌ OCR failed on page {i + 1}: {e}")

#                 if all_text:
#                     combined_text = "\n".join(all_text)
#                     st.write(
#                         f"✅ **OCR Complete!** Extracted {len(combined_text)} characters"
#                     )

#                     # Store result
#                     st.session_state["ocr_text"] = combined_text
#                     st.session_state["ocr_file"] = ocr_pdf_path

#                     status.update(label="✅ OCR Processing Complete!", state="complete")
#                 else:
#                     st.error("❌ No text could be extracted")

#         except Exception as e:
#             st.error(f"❌ **OCR Error**: {str(e)}")
#             import traceback

#             st.code(traceback.format_exc())
#     else:
#         st.error("❌ PDF file not found!")

# # Show results if available
# if "ocr_text" in st.session_state and st.session_state["ocr_text"]:
#     st.write("---")
#     st.write("## 📄 OCR Results")

#     ocr_text = st.session_state["ocr_text"]
#     ocr_file = st.session_state["ocr_file"]

#     st.write(f"**File**: {os.path.basename(ocr_file)}")
#     st.write(f"**Extracted**: {len(ocr_text)} characters")

#     # Show preview
#     st.text_area(
#         "OCR Text Preview", ocr_text[:1000], height=200, key="ocr_result_preview"
#     )

#     # Add to vector store button
#     if st.button("Add OCR Text to Vector Store", key="add_ocr_final"):
#         try:
#             # Create document
#             doc = Document(
#                 page_content=ocr_text,
#                 metadata={
#                     "source": ocr_file,
#                     "file_name": os.path.basename(ocr_file),
#                     "file_type": "pdf",
#                     "extraction_method": "ocr",
#                     "text_length": len(ocr_text),
#                 },
#             )

#             # Create chunks
#             text_splitter = CharacterTextSplitter(
#                 separator="\n",
#                 chunk_size=1000,
#                 chunk_overlap=200,
#             )
#             chunks = text_splitter.split_documents([doc])

#             # Add to vector store
#             if "vectorstore" in st.session_state:
#                 st.session_state.vectorstore.add_documents(chunks)
#                 save_vectorstore(st.session_state.vectorstore)
#                 save_local_documents([doc], [ocr_file], "pdf")

#                 st.success(
#                     f"✅ **Success!** Added {len(chunks)} chunks to vector store"
#                 )

#                 st.success(f"✅ **Added {len(chunks)} chunks to vector store!**")

#                 # Clear results without immediate rerun
#                 del st.session_state["ocr_text"]
#                 del st.session_state["ocr_file"]

#                 # Show instruction instead of auto-rerun
#                 st.info("🔄 **Refresh the page to clear this section**")
#             else:
#                 st.error("❌ No vector store found!")

#         except Exception as e:
#             st.error(f"❌ Error adding to vector store: {str(e)}")

#     # Clear results button
#     if st.button("Clear OCR Results", key="clear_ocr_simple"):
#         del st.session_state["ocr_text"]
#         del st.session_state["ocr_file"]
#         st.rerun()


# Show loaded documents/URLs in sidebar
# with st.sidebar:
# st.subheader("Processed Content")
# if st.session_state.processed_urls:
# Get content statistics
# stats = get_content_statistics()
# st.write(f"**Total items:** {len(st.session_state.processed_urls)}")
# # Show breakdown by type
# if stats.get("web_pages", 0) > 0:
#     st.write(f"📄 Web pages: {stats['web_pages']}")
# if stats.get("html_files", 0) > 0:
#     st.write(f"📝 HTML files: {stats['html_files']}")
# if stats.get("docx_files", 0) > 0:
#     st.write(f"📘 DOCX files: {stats['docx_files']}")
# if stats.get("pdf_files", 0) > 0:
#     st.write(f"📕 PDF files: {stats['pdf_files']}")
# if stats.get("pptx_files", 0) > 0:
#     st.write(f"📊 PowerPoint files: {stats['pptx_files']}")
# if stats.get("xlsx_files", 0) > 0:
#     st.write(f"📈 Excel files: {stats['xlsx_files']}")
# with st.expander("View all items", expanded=False):
#     for item in st.session_state.processed_urls:
#         if item.startswith("http"):
#             st.write(f"🌐 {item}")
#         else:
#             st.write(f"📁 {os.path.basename(item)}")

#     # Initialize confirmation state
#     if "db_delete_confirmation" not in st.session_state:
#         st.session_state.db_delete_confirmation = False

#     col1, col2 = st.columns(2)
#     with col1:
#         # Display different buttons based on confirmation state
#         if not st.session_state.db_delete_confirmation:
#             # Initial button
#             if st.button("Clear Database", type="primary"):
#                 st.session_state.db_delete_confirmation = True
#                 st.rerun()
#         else:
#             # Show warning and confirmation buttons
#             st.warning(
#                 "⚠️ Are you sure you want to delete all data? This cannot be undone."
#             )
#             # Yes button
#             if st.button("Yes, Delete", type="primary"):
#                 # Clear tables but keep structure
#                 conn = get_db_connection()
#                 cursor = conn.cursor()
#                 # Delete in the correct order to respect foreign key constraints
#                 cursor.execute("DELETE FROM pages")
#                 cursor.execute("DELETE FROM websites")
#                 cursor.execute("DELETE FROM documents")
#                 conn.commit()
#                 # Remove vector store files
#                 import shutil

#     if os.path.exists(os.path.join(DATA_DIR, "faiss_index")):
#         shutil.rmtree(os.path.join(DATA_DIR, "faiss_index"))
#     # Clear session state but keep chat history
#     for key in [
#         "vectorstore",
#         "conversation_chain",
#         "processed_urls",
#         "db_delete_confirmation",
#     ]:
#         if key in st.session_state:
#             del st.session_state[key]
#     st.success("Database cleared successfully!")
#     st.rerun()
# # No button
# if st.button("No, Cancel"):
#     st.session_state.db_delete_confirmation = False
#     st.rerun()

# with col2:
# if st.button("Clear Chat History", type="secondary"):
# if st.button("Clear Chat History"):
#             # Clear only chat history
#             conn = get_db_connection()
#             cursor = conn.cursor()
#             cursor.execute("DELETE FROM chat_history")
#             conn.commit()
#             # Clear chat history from session state
#             st.session_state.chat_history = []
#             # If conversation chain exists, reset its memory
#             if "conversation_chain" in st.session_state:
#                 # Create a new memory object
#                 llm = ChatGroq(model="llama-3.3-70b-versatile", temperature=0)
#                 memory = ConversationBufferMemory(
#                     llm=llm,
#                     output_key="answer",
#                     memory_key="chat_history",
#                     return_messages=True,
#                 )
#                 # Update the chain with the new memory
#                 st.session_state.conversation_chain.memory = memory
#             st.success("Chat history cleared!")
#             st.rerun()
# else:
#   st.info("No content processed yet. Enter paths or URLs above to begin.")


# Add this after the sidebar section and before the chat interface
# if "vectorstore" in st.session_state:
#     st.divider()
#     export_container = st.container()
#     with export_container:
#         st.subheader("Data Export")
#         # Create two columns for export buttons
#         col1, col2 = st.columns(2)
#         # Export Vector Database
#         with col1:
#             # Add export button for FAISS vectors
#             if st.button("Export Vectors to CSV", type="primary"):
#                 with st.spinner("Exporting vectors..."):
#                     # Define the export function
#                     def export_faiss_vectors(filename="faiss_vectors.csv"):
#                         """Export FAISS vectors to CSV for external visualization"""
#                         try:
#                             vs = st.session_state.vectorstore
#                             # Get all the vectors and their IDs
#                             vectors = []
#                             ids = []
#                             metadata = []
#                             # Different vectorstores have different structures
#                             if hasattr(vs, "index_to_docstore_id") and hasattr(
#                                 vs, "index"
#                             ):
#                                 for i in range(len(vs.index_to_docstore_id)):
#                                     doc_id = vs.index_to_docstore_id.get(i)
#                                     if doc_id and doc_id in vs.docstore._dict:
#                                         # Get the document
#                                         doc = vs.docstore._dict[doc_id]
#                                         # Get the vector
#                                         vector = vs.index.reconstruct(i)
#                                         vectors.append(vector)
#                                         ids.append(doc_id)
#                                         # Enhanced metadata for different document types
#                                         meta = {
#                                             "source": doc.metadata.get("source", ""),
#                                             "content_preview": doc.page_content[:100],
#                                         }
#                                         if doc.metadata.get("type") == "web":
#                                             meta["title"] = doc.metadata.get(
#                                                 "title", ""
#                                             )
#                                             meta["url"] = doc.metadata.get("url", "")
#                                         else:
#                                             meta["file_name"] = doc.metadata.get(
#                                                 "file_name", ""
#                                             )
#                                             meta["file_type"] = doc.metadata.get(
#                                                 "file_type", ""
#                                             )
#                                         metadata.append(meta)
#                             # Export to CSV
#                             import pandas as pd
#                             import numpy as np

#                             # Create a dataframe with metadata
#                             meta_df = pd.DataFrame(metadata)
#                             # Create a dataframe with vectors
#                             vector_df = pd.DataFrame(np.array(vectors))
#                             # Combine the dataframes
#                             result_df = pd.concat([meta_df, vector_df], axis=1)
#                             # Save to CSV
#                             result_df.to_csv(filename, index=False)
#                             return (
#                                 True,
#                                 f"Exported {len(vectors)} vectors to {filename}",
#                             )
#                         except Exception as e:
#                             return False, f"Error exporting vectors: {str(e)}"

#                     # Export the vectors to a file in the data directory
#                     export_path = os.path.join(DATA_DIR, "faiss_vectors.csv")
#                     success, message = export_faiss_vectors(export_path)
#                     if success:
#                         st.success(message)
#                         # Create a download link for the exported file
#                         with open(export_path, "rb") as file:
#                             st.download_button(
#                                 label="Download Vector CSV",
#                                 data=file,
#                                 file_name="faiss_vectors.csv",
#                                 mime="text/csv",
#                             )
#                     else:
#                         st.error(message)

#         # Export SQLite Database
#         with col2:
#             if st.button("Export SQLite Database", type="primary"):
#                 with st.spinner("Preparing SQLite database for export..."):
#                     try:
#                         # Create a copy of the database for export
#                         import shutil

#                         export_db_path = os.path.join(
#                             DATA_DIR, "website_data_export.db"
#                         )
#                         # Close any existing connections before copying
#                         close_db_connection()
#                         # Copy the database file
#                         shutil.copy2(DB_PATH, export_db_path)
#                         st.success("SQLite database prepared for download!")
#                         # Create download button for the SQLite file
#                         with open(export_db_path, "rb") as file:
#                             st.download_button(
#                                 label="Download SQLite Database",
#                                 data=file,
#                                 file_name="website_data.db",
#                                 mime="application/x-sqlite3",
#                             )
#                     except Exception as e:
#                         st.error(f"Error exporting SQLite database: {str(e)}")


# Initialize conversation chain if vectorstore exists
if "vectorstore" in st.session_state and "conversation_chain" not in st.session_state:
    st.session_state.conversation_chain = create_chain_with_custom_retriever(
        st.session_state.vectorstore,
        k=50,  # Retrieve top 50 chunks
    )

###################### Chat Handler (Begin) ####################
# Chat interface
if "vectorstore" in st.session_state:
    st.divider()
    st.subheader("Ask Claudia:")
    # Display chat history
    for message in st.session_state.chat_history:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

# Chat input
user_input = st.chat_input("Ask for Office Information...")

if user_input and user_input.strip():
    # Add user message to chat history
    st.session_state.chat_history.append({"role": "user", "content": user_input})
    with st.chat_message("user"):
        st.write("Lowly Human")
        st.markdown(user_input)

    # Generate assistant response
    with st.chat_message("assistant"):
        st.write("Claudia")
        with st.spinner("Thinking..."):
            assistant_response = handle_improved_query_routing(
                user_input, st.session_state.vectorstore
            )
            st.markdown(assistant_response)

    # Add assistant response to chat history
    st.session_state.chat_history.append(
        {"role": "assistant", "content": assistant_response}
    )

    # Save chat history to database
    save_chat_history(st.session_state.chat_history)


##################################### Chat Handler (End)
if st.button("Clear Chat History"):
    # Clear only chat history
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("DELETE FROM chat_history")
    conn.commit()
    # Clear chat history from session state
    st.session_state.chat_history = []
    # If conversation chain exists, reset its memory
    if "conversation_chain" in st.session_state:
        # Create a new memory object
        llm = ChatGroq(model="llama-3.3-70b-versatile", temperature=0)
        memory = ConversationBufferMemory(
            llm=llm,
            output_key="answer",
            memory_key="chat_history",
            return_messages=True,
        )
        # Update the chain with the new memory
        st.session_state.conversation_chain.memory = memory
    st.success("Chat history cleared!")
    st.rerun()
