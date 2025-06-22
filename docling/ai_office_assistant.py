# Add these imports with the existing ones
from pptx import Presentation
import pandas as pd
import openpyxl
import os
from dotenv import load_dotenv
import streamlit as st
from docling.document_converter import DocumentConverter
from langchain.schema import Document
from langchain_text_splitters.character import CharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_groq import ChatGroq
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
import xml.etree.ElementTree as ET
from typing import List
from urllib.parse import urljoin
import requests
import time
import random
import sqlite3
import datetime
import threading
from pathlib import Path
import glob
import re

# Load the environment variables
load_dotenv()
working_dir = os.path.dirname(os.path.abspath(__file__))
# Define data directory for persistent storage
DATA_DIR = os.path.join(working_dir, "data")
os.makedirs(DATA_DIR, exist_ok=True)
DB_PATH = os.path.join(DATA_DIR, "website_data.db")
# Thread-local storage for database connections
_local = threading.local()


def get_all_names_from_vectorstore_improved(vectorstore):
    """Extract ONLY real names from the vectorstore - FIXED for real names with punctuation"""
    try:
        all_names = set()

        if hasattr(vectorstore, "docstore") and hasattr(
            vectorstore, "index_to_docstore_id"
        ):
            for i in range(len(vectorstore.index_to_docstore_id)):
                doc_id = vectorstore.index_to_docstore_id.get(i)
                if doc_id and doc_id in vectorstore.docstore._dict:
                    doc = vectorstore.docstore._dict[doc_id]

                    # Only process Excel chunks
                    if doc.metadata.get("file_type") in ["xlsx", "xls"]:
                        content = doc.page_content

                        # ONLY use the most precise pattern for your data structure
                        name_pattern = (
                            r"FIRST_NAME:\s*([^|]+?)\s*\|\s*LAST_NAME:\s*([^|]+?)\s*\|"
                        )
                        matches = re.findall(name_pattern, content, re.IGNORECASE)

                        for first, last in matches:
                            first = first.strip()
                            last = last.strip()

                            # More realistic filtering for real names
                            if (
                                first
                                and last
                                and len(first) >= 2
                                and len(last) >= 2
                                and
                                # Allow letters, spaces, periods, commas, apostrophes, hyphens
                                re.match(r"^[A-Za-z\s\.\,\'\-]+$", first)
                                and re.match(
                                    r"^[A-Za-z\s\.\,\'\-IVX]+$", last
                                )  # Allow Roman numerals
                                and
                                # Exclude obvious non-names
                                not any(
                                    word.lower() in first.lower()
                                    for word in [
                                        "record",
                                        "structured",
                                        "data",
                                        "sheet",
                                        "column",
                                        "total",
                                        "fee",
                                    ]
                                )
                                and not any(
                                    word.lower() in last.lower()
                                    for word in [
                                        "record",
                                        "structured",
                                        "data",
                                        "sheet",
                                        "column",
                                        "total",
                                        "fee",
                                    ]
                                )
                                and
                                # Must start with a letter
                                first[0].isalpha()
                                and last[0].isalpha()
                            ):
                                full_name = f"{first} {last}"
                                all_names.add(full_name)

        return sorted(list(all_names))

    except Exception as e:
        st.error(f"Error extracting names: {str(e)}")
        return []


def get_db_connection():
    """Get a thread-specific database connection"""
    if not hasattr(_local, "conn"):
        _local.conn = sqlite3.connect(DB_PATH, check_same_thread=False)
        # Enable foreign keys
        _local.conn.execute("PRAGMA foreign_keys = ON")
    return _local.conn


def close_db_connection():
    """Close the thread-specific database connection if it exists"""
    if hasattr(_local, "conn"):
        _local.conn.close()
        delattr(_local, "conn")


def setup_database():
    """Set up SQLite database for document storage"""
    conn = get_db_connection()
    cursor = conn.cursor()
    # Create enhanced tables for all document types
    cursor.execute("""
    CREATE TABLE IF NOT EXISTS websites (
        id INTEGER PRIMARY KEY,
        sitemap_url TEXT UNIQUE,
        processed_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    """)
    cursor.execute("""
    CREATE TABLE IF NOT EXISTS pages (
        id INTEGER PRIMARY KEY,
        website_id INTEGER,
        url TEXT UNIQUE,
        title TEXT,
        content TEXT,
        processed_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (website_id) REFERENCES websites(id)
    )
    """)
    # New table for local documents
    cursor.execute("""
    CREATE TABLE IF NOT EXISTS documents (
        id INTEGER PRIMARY KEY,
        file_path TEXT UNIQUE,
        file_name TEXT,
        file_type TEXT,
        source_type TEXT,
        content TEXT,
        processed_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    """)
    # Create table for chat history
    cursor.execute("""
    CREATE TABLE IF NOT EXISTS chat_history (
        id INTEGER PRIMARY KEY,
        role TEXT,
        content TEXT,
        timestamp TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    """)
    conn.commit()
    return True


def save_vectorstore(vectorstore, save_path=DATA_DIR):
    """Save FAISS vectorstore to disk"""
    os.makedirs(save_path, exist_ok=True)
    vectorstore.save_local(f"{save_path}/faiss_index")


def load_vectorstore(load_path=DATA_DIR):
    """Load FAISS vectorstore from disk"""
    if os.path.exists(f"{load_path}/faiss_index"):
        embeddings = HuggingFaceEmbeddings()
        vectorstore = FAISS.load_local(
            f"{load_path}/faiss_index",
            embeddings,
            allow_dangerous_deserialization=True,
        )
        return vectorstore
    return None


def save_processed_content(sitemap_url, processed_urls, documents):
    """Save processed content to SQLite database"""
    conn = get_db_connection()
    cursor = conn.cursor()
    try:
        # Insert or update website record
        cursor.execute(
            "INSERT OR REPLACE INTO websites (sitemap_url, processed_date) VALUES (?, ?)",
            (sitemap_url, datetime.datetime.now()),
        )
        website_id = (
            cursor.lastrowid
            or cursor.execute(
                "SELECT id FROM websites WHERE sitemap_url = ?", (sitemap_url,)
            ).fetchone()[0]
        )
        # Insert document records
        for doc, url in zip(documents, processed_urls):
            # Extract title if available in metadata
            title = doc.metadata.get("title", os.path.basename(url))
            cursor.execute(
                "INSERT OR REPLACE INTO pages (website_id, url, title, content, processed_date) VALUES (?, ?, ?, ?, ?)",
                (website_id, url, title, doc.page_content, datetime.datetime.now()),
            )
        conn.commit()
    except Exception as e:
        conn.rollback()
        st.error(f"Database error: {str(e)}")


def save_local_documents(documents, file_paths, file_type):
    """Save locally processed documents to SQLite database"""
    conn = get_db_connection()
    cursor = conn.cursor()
    try:
        for doc, file_path in zip(documents, file_paths):
            file_name = os.path.basename(file_path)
            cursor.execute(
                "INSERT OR REPLACE INTO documents (file_path, file_name, file_type, source_type, content, processed_date) VALUES (?, ?, ?, ?, ?, ?)",
                (
                    file_path,
                    file_name,
                    file_type,
                    "local_directory",
                    doc.page_content,
                    datetime.datetime.now(),
                ),
            )
        conn.commit()
    except Exception as e:
        conn.rollback()
        st.error(f"Database error: {str(e)}")


def load_processed_content(sitemap_url=None):
    """Load processed content from SQLite database"""
    conn = get_db_connection()
    cursor = conn.cursor()
    documents = []
    processed_urls = []
    # Load web pages
    if sitemap_url:
        cursor.execute(
            """
            SELECT p.url, p.title, p.content
            FROM pages p
            JOIN websites w ON p.website_id = w.id
            WHERE w.sitemap_url = ?
            """,
            (sitemap_url,),
        )
    else:
        cursor.execute("SELECT url, title, content FROM pages")
    results = cursor.fetchall()
    for url, title, content in results:
        doc = Document(
            page_content=content,
            metadata={"source": url, "url": url, "title": title, "type": "web"},
        )
        documents.append(doc)
        processed_urls.append(url)
    # Load local documents
    cursor.execute("SELECT file_path, file_name, content, file_type FROM documents")
    results = cursor.fetchall()
    for file_path, file_name, content, file_type in results:
        doc = Document(
            page_content=content,
            metadata={
                "source": file_path,
                "file_name": file_name,
                "file_type": file_type,
                "type": "local",
            },
        )
        documents.append(doc)
        processed_urls.append(file_path)
    return documents, processed_urls


def save_chat_history(chat_history):
    """Save chat history to database"""
    conn = get_db_connection()
    cursor = conn.cursor()
    try:
        # Clear existing history
        cursor.execute("DELETE FROM chat_history")
        # Insert new history
        for idx, message in enumerate(chat_history):
            cursor.execute(
                "INSERT INTO chat_history (id, role, content) VALUES (?, ?, ?)",
                (idx, message["role"], message["content"]),
            )
        conn.commit()
    except Exception as e:
        conn.rollback()
        st.error(f"Error saving chat history: {str(e)}")


def load_chat_history():
    """Load chat history from database"""
    conn = get_db_connection()
    cursor = conn.cursor()
    try:
        cursor.execute("SELECT role, content FROM chat_history ORDER BY id")
        results = cursor.fetchall()
        chat_history = []
        for role, content in results:
            chat_history.append({"role": role, "content": content})
        return chat_history
    except:
        return []


def process_powerpoint_file(file_path: str) -> Document:
    """Process a PowerPoint file and extract text content"""
    try:
        # First try with Docling (it might support PPTX)
        try:
            converter = DocumentConverter()
            result = converter.convert(file_path)
            title = os.path.basename(file_path)
            # Try to get markdown format first
            try:
                full_text = result.document.export_to_markdown()
                if full_text and full_text.strip():
                    return Document(
                        page_content=full_text,
                        metadata={
                            "source": file_path,
                            "file_name": title,
                            "file_type": "pptx",
                            "format": "markdown",
                        },
                    )
            except AttributeError:
                pass
            # Try other content extraction methods
            if hasattr(result.document, "text"):
                text_content = result.document.text
            elif hasattr(result.document, "get_text"):
                text_content = result.document.get_text()
            else:
                text_content = str(result.document)
            if text_content and text_content.strip():
                return Document(
                    page_content=text_content,
                    metadata={
                        "source": file_path,
                        "file_name": title,
                        "file_type": "pptx",
                    },
                )
        except Exception:
            # If Docling fails, fall back to python-pptx
            pass
        # Fallback to manual extraction using python-pptx
        prs = Presentation(file_path)
        text_content = []
        slide_count = 0
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            slide_text.append(f"\n--- Slide {slide_num} ---\n")
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                # Handle tables in slides
                if shape.shape_type == 19:  # Table
                    try:
                        table_text = []
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                row_text.append(cell.text.strip())
                            table_text.append(" | ".join(row_text))
                        slide_text.append("\n".join(table_text))
                    except:
                        pass
            if len(slide_text) > 1:  # More than just the slide header
                text_content.extend(slide_text)
                slide_count += 1
        if text_content:
            full_text = "\n".join(text_content)
            title = os.path.basename(file_path)
            return Document(
                page_content=full_text,
                metadata={
                    "source": file_path,
                    "file_name": title,
                    "file_type": "pptx",
                    "slide_count": slide_count,
                    "extraction_method": "python-pptx",
                },
            )
        return None
    except Exception as e:
        st.warning(f"Error processing PowerPoint file {file_path}: {str(e)}")
        return None


def process_excel_file(file_path: str) -> Document:
    """Process an Excel file with better structure preservation"""
    try:
        text_content = []
        sheet_count = 0
        try:
            excel_file = pd.ExcelFile(file_path, engine="openpyxl")
        except:
            try:
                excel_file = pd.ExcelFile(file_path, engine="xlrd")
            except:
                excel_file = pd.ExcelFile(file_path)
        for sheet_name in excel_file.sheet_names:
            try:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                if df.empty:
                    continue
                sheet_text = [f"\n=== SHEET: {sheet_name} ===\n"]
                # Add structured data representation
                df_clean = df.fillna("")
                # For name-related data, create a more structured format
                if any(
                    col.lower() in ["name", "first_name", "last_name", "full_name"]
                    for col in df_clean.columns
                ):
                    sheet_text.append("STRUCTURED DATA - PEOPLE RECORDS:")
                    sheet_text.append("-" * 40)
                    for idx, row in df_clean.iterrows():
                        row_parts = []
                        for col, val in row.items():
                            if str(val).strip():
                                row_parts.append(f"{col}: {val}")
                        if row_parts:
                            sheet_text.append(
                                f"Record {idx + 1}: {' | '.join(row_parts)}"
                            )
                else:
                    # Standard table format for non-name data
                    if not df_clean.columns.empty:
                        headers = " | ".join(str(col) for col in df_clean.columns)
                        sheet_text.append(f"COLUMNS: {headers}")
                        sheet_text.append("-" * len(headers))
                    for idx, row in df_clean.iterrows():
                        row_text = " | ".join(
                            str(val) for val in row.values if str(val).strip()
                        )
                        if row_text.strip():
                            sheet_text.append(f"Row {idx + 1}: {row_text}")
                sheet_text.append(
                    f"\n[SHEET SUMMARY: {len(df)} total rows, {len(df.columns)} columns]"
                )
                sheet_text.append("=" * 50)
                text_content.extend(sheet_text)
                sheet_count += 1
            except Exception as e:
                st.warning(f"Error reading sheet '{sheet_name}': {str(e)}")
                continue
        if text_content:
            full_text = "\n".join(text_content)
            title = os.path.basename(file_path)
            return Document(
                page_content=full_text,
                metadata={
                    "source": file_path,
                    "file_name": title,
                    "file_type": "xlsx",
                    "sheet_count": sheet_count,
                    "extraction_method": "structured_pandas",
                    "total_content_length": len(full_text),
                },
            )
        return None
    except Exception as e:
        st.warning(f"Error processing Excel file {file_path}: {str(e)}")
        return None


def process_document_file(file_path: str, file_type: str) -> Document:
    """Process a single document file using appropriate method with OCR fallback for PDFs"""
    try:
        # Handle different file types with specialized functions
        if file_type.lower() == "pptx":
            return process_powerpoint_file(file_path)
        elif file_type.lower() in ["xlsx", "xls"]:
            return process_excel_file(file_path)
        elif file_type.lower() == "pdf":
            # For PDFs, try Docling first, then OCR if needed
            return process_pdf_with_intelligent_fallback(file_path)
        else:
            # Use original Docling approach for other formats
            converter = DocumentConverter()
            result = converter.convert(file_path)
            title = os.path.basename(file_path)

            # Try to get the full text as markdown
            try:
                full_text = result.document.export_to_markdown()
                if full_text and full_text.strip():
                    return Document(
                        page_content=full_text,
                        metadata={
                            "source": file_path,
                            "file_name": title,
                            "file_type": file_type,
                            "format": "markdown",
                        },
                    )
            except AttributeError:
                pass

            # If no markdown export, try to access content directly
            if hasattr(result.document, "text"):
                text_content = result.document.text
            elif hasattr(result.document, "get_text"):
                text_content = result.document.get_text()
            elif hasattr(result, "text"):
                text_content = result.text
            else:
                text_content = str(result.document)

            if text_content and text_content.strip():
                return Document(
                    page_content=text_content,
                    metadata={
                        "source": file_path,
                        "file_name": title,
                        "file_type": file_type,
                    },
                )
            return None
    except Exception as e:
        st.warning(f"Error processing file {file_path}: {str(e)}")
        return None


def process_pdf_with_intelligent_fallback(file_path: str) -> Document:
    """Process PDF with Docling first, then OCR fallback if needed"""
    try:
        st.write(f"**Processing PDF**: {os.path.basename(file_path)}")

        # First try regular Docling extraction
        try:
            converter = DocumentConverter()
            result = converter.convert(file_path)

            # Try to get markdown format first
            try:
                full_text = result.document.export_to_markdown()
                if full_text and full_text.strip() and len(full_text) > 100:
                    # Check if we got meaningful content (not just image placeholders)
                    if (
                        full_text.count("<!-- image -->") / len(full_text) < 0.1
                    ):  # Less than 10% image tags
                        st.write("✅ **Regular text extraction successful**")
                        return Document(
                            page_content=full_text,
                            metadata={
                                "source": file_path,
                                "file_name": os.path.basename(file_path),
                                "file_type": "pdf",
                                "format": "markdown",
                                "extraction_method": "docling",
                            },
                        )
            except AttributeError:
                pass

            # Try other content extraction methods
            if hasattr(result.document, "text"):
                text_content = result.document.text
            elif hasattr(result.document, "get_text"):
                text_content = result.document.get_text()
            else:
                text_content = str(result.document)

            # Check if we got meaningful text (more than just image tags)
            if (
                text_content
                and len(text_content) > 100
                and "<!-- image -->" not in text_content
            ):
                st.write("✅ **Regular text extraction successful**")
                return Document(
                    page_content=text_content,
                    metadata={
                        "source": file_path,
                        "file_name": os.path.basename(file_path),
                        "file_type": "pdf",
                        "extraction_method": "docling",
                    },
                )
        except Exception as e:
            st.write(f"⚠️ **Regular extraction failed**: {str(e)}")

        # If regular extraction failed or gave poor results, try OCR
        st.write("🔍 **Falling back to OCR extraction...**")
        return process_pdf_with_ocr_internal(file_path)

    except Exception as e:
        st.warning(f"Error processing PDF {file_path}: {str(e)}")
        return None


def process_pdf_with_ocr_internal(file_path: str) -> Document:
    """Internal OCR processing function (extracted from your existing code)"""
    try:
        import pytesseract
        from pdf2image import convert_from_path
        from PIL import Image
        import tempfile
        import os

        st.write("🔍 **Starting OCR extraction...**")

        # Convert PDF pages to images
        with tempfile.TemporaryDirectory() as temp_dir:
            try:
                # Convert PDF to images
                pages = convert_from_path(file_path, dpi=300)
                st.write(f"📄 **Converted to {len(pages)} images**")

                extracted_text = []
                for i, page in enumerate(pages):
                    st.write(f"🔍 **Processing page {i + 1}/{len(pages)}...**")
                    # Save page as temporary image
                    temp_image_path = os.path.join(temp_dir, f"page_{i + 1}.png")
                    page.save(temp_image_path, "PNG")

                    # Extract text using OCR
                    try:
                        page_text = pytesseract.image_to_string(page)
                        if page_text.strip():
                            extracted_text.append(f"\n--- Page {i + 1} ---\n")
                            extracted_text.append(page_text.strip())
                        else:
                            st.write(f"⚠️ **No text found on page {i + 1}**")
                    except Exception as e:
                        st.write(f"❌ **OCR failed for page {i + 1}**: {str(e)}")
                        continue

                if extracted_text:
                    full_text = "\n".join(extracted_text)
                    st.write(
                        f"✅ **OCR successful! Extracted {len(full_text)} characters**"
                    )

                    return Document(
                        page_content=full_text,
                        metadata={
                            "source": file_path,
                            "file_name": os.path.basename(file_path),
                            "file_type": "pdf",
                            "extraction_method": "ocr",
                            "pages_processed": len(pages),
                            "text_length": len(full_text),
                        },
                    )
                else:
                    st.error("❌ **No text could be extracted via OCR**")
                    return None

            except Exception as e:
                st.error(f"❌ **PDF to image conversion failed**: {str(e)}")
                return None

    except ImportError:
        st.error(
            "❌ **OCR libraries not installed**. Please install: pip install pytesseract pillow pdf2image"
        )
        return None
    except Exception as e:
        st.error(f"❌ **OCR processing failed**: {str(e)}")
        return None


def process_local_directory(
    directory_path: str, file_extension: str, max_files: int = None
):
    """Process all files of a specific type in a directory"""
    if not os.path.exists(directory_path):
        st.error(f"Directory not found: {directory_path}")
        return [], []
    # Get all files with the specified extension
    pattern = os.path.join(directory_path, f"*.{file_extension}")
    file_paths = glob.glob(pattern)
    if not file_paths:
        st.warning(f"No .{file_extension} files found in {directory_path}")
        return [], []
    st.success(f"Found {len(file_paths)} .{file_extension} files")
    # Limit number of files if specified
    if max_files and max_files < len(file_paths):
        st.info(f"Processing only {max_files} files out of {len(file_paths)}")
        file_paths = file_paths[:max_files]
    all_documents = []
    processed_files = []
    # Process each file
    progress_bar = st.progress(0)
    status_text = st.empty()
    for idx, file_path in enumerate(file_paths):
        status_text.text(
            f"Processing file {idx + 1}/{len(file_paths)}: {os.path.basename(file_path)}"
        )
        try:
            document = process_document_file(file_path, file_extension)
            if document:
                all_documents.append(document)
                processed_files.append(file_path)
            progress_bar.progress((idx + 1) / len(file_paths))
            # Add small delay to prevent overwhelming the system
            time.sleep(0.1)
        except Exception as e:
            st.warning(f"Failed to process {file_path}: {str(e)}")
            continue
    progress_bar.empty()
    status_text.empty()
    if not all_documents:
        st.error(f"No documents could be extracted from .{file_extension} files")
    return all_documents, processed_files


def get_sitemap_urls(sitemap_url: str) -> List[str]:
    """
    Extracts all URLs from a sitemap XML file.
    Handles both regular sitemaps and sitemap indexes.
    """
    try:
        response = requests.get(sitemap_url, timeout=30)
        response.raise_for_status()
        root = ET.fromstring(response.content)
        # Define XML namespaces
        namespaces = {
            "sm": "http://www.sitemaps.org/schemas/sitemap/0.9",
            "xhtml": "http://www.w3.org/1999/xhtml",
        }
        urls = []
        # Check if this is a sitemap index (contains other sitemaps)
        sitemap_tags = root.findall(".//sm:sitemap", namespaces)
        if sitemap_tags:
            # This is a sitemap index
            for sitemap_tag in sitemap_tags:
                loc_tag = sitemap_tag.find("sm:loc", namespaces)
                if loc_tag is not None and loc_tag.text:
                    # Recursively get URLs from each child sitemap
                    child_sitemap_url = loc_tag.text.strip()
                    child_urls = get_sitemap_urls(child_sitemap_url)
                    urls.extend(child_urls)
        else:
            # This is a regular sitemap
            url_tags = root.findall(".//sm:url", namespaces)
            for url_tag in url_tags:
                loc_tag = url_tag.find("sm:loc", namespaces)
                if loc_tag is not None and loc_tag.text:
                    urls.append(loc_tag.text.strip())
        return urls
    except Exception as e:
        st.error(f"Error processing sitemap: {str(e)}")
        return []


def fetch_html_content(url: str) -> str:
    """Fetch HTML content from a URL"""
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
        }
        response = requests.get(url, headers=headers, timeout=30)
        response.raise_for_status()
        return response.text
    except Exception as e:
        st.warning(f"Error fetching HTML from {url}: {str(e)}")
        return ""


def process_url_with_docling(url: str, temp_dir: str) -> Document:
    """Process a URL's HTML content with Docling"""
    html_content = fetch_html_content(url)
    if not html_content:
        return None
    # Create a temporary HTML file
    url_filename = url.replace(":", "_").replace("/", "_").replace(".", "_") + ".html"
    temp_file_path = os.path.join(temp_dir, url_filename)
    with open(temp_file_path, "w", encoding="utf-8") as f:
        f.write(html_content)
    # Process with Docling
    try:
        converter = DocumentConverter()
        result = converter.convert(temp_file_path)
        # Extract page title if possible
        title = url
        try:
            title_match = re.search("<title>(.*?)</title>", html_content, re.IGNORECASE)
            if title_match:
                title = title_match.group(1)
        except:
            pass
            # Try to get the full text as markdown
        try:
            full_text = result.document.export_to_markdown()
            if full_text and full_text.strip():
                return Document(
                    page_content=full_text,
                    metadata={
                        "source": url,
                        "url": url,
                        "title": title,
                        "format": "markdown",
                    },
                )
        except AttributeError:
            pass
        # If no markdown export, try to access content directly
        if hasattr(result.document, "text"):
            text_content = result.document.text
        elif hasattr(result.document, "get_text"):
            text_content = result.document.get_text()
        elif hasattr(result, "text"):
            text_content = result.text
        else:
            text_content = str(result.document)
        if text_content and text_content.strip():
            return Document(
                page_content=text_content,
                metadata={
                    "source": url,
                    "url": url,
                    "title": title,
                },
            )
        # If still no content, try to access raw content
        if hasattr(result, "content"):
            content = result.content
            if content and str(content).strip():
                return Document(
                    page_content=str(content),
                    metadata={
                        "source": url,
                        "url": url,
                        "title": title,
                    },
                )
        return None
    except Exception as e:
        st.warning(f"Error processing URL with Docling: {url} - {str(e)}")
        return None
    finally:
        # Clean up the temp file
        if os.path.exists(temp_file_path):
            try:
                os.remove(temp_file_path)
            except:
                pass


def process_sitemap(sitemap_url: str, max_urls: int = None):
    """Process a sitemap and extract content from all URLs"""
    with st.spinner("Fetching sitemap URLs..."):
        urls = get_sitemap_urls(sitemap_url)
        if not urls:
            st.error("No URLs found in the sitemap.")
            return [], []
        st.success(f"Found {len(urls)} URLs in the sitemap.")
        if max_urls and max_urls < len(urls):
            st.info(f"Processing only {max_urls} URLs out of {len(urls)}.")
            urls = urls[:max_urls]
    # Create a temporary directory for HTML files
    temp_dir = os.path.join(working_dir, "temp_html")
    os.makedirs(temp_dir, exist_ok=True)
    all_documents = []
    processed_urls = []
    # Process each URL
    progress_bar = st.progress(0)
    status_text = st.empty()
    for idx, url in enumerate(urls):
        status_text.text(f"Processing URL {idx + 1}/{len(urls)}: {url}")
        try:
            # Add some delay to avoid overloading the server
            time.sleep(random.uniform(0.5, 2.0))
            document = process_url_with_docling(url, temp_dir)
            if document:
                all_documents.append(document)
                processed_urls.append(url)
            progress_bar.progress((idx + 1) / len(urls))
        except Exception as e:
            st.warning(f"Failed to process {url}: {str(e)}")
            continue
    progress_bar.empty()
    status_text.empty()
    # Clean up the temp directory
    try:
        import shutil

        shutil.rmtree(temp_dir)
    except:
        pass
    if not all_documents:
        st.error("No documents could be extracted from the sitemap URLs.")
    return all_documents, processed_urls


def setup_vectorstore(documents):
    embeddings = HuggingFaceEmbeddings()
    # Different splitting strategy for different document types
    doc_chunks = []
    for doc in documents:
        if doc.metadata.get("file_type") in ["xlsx", "xls"]:
            # For Excel files, use larger chunks to preserve record integrity
            text_splitter = CharacterTextSplitter(
                separator="\n",
                chunk_size=2000,  # Larger chunks for tabular data
                chunk_overlap=100,  # Smaller overlap
            )
        else:
            # Standard splitting for other documents
            text_splitter = CharacterTextSplitter(
                separator="\n",
                chunk_size=1000,
                chunk_overlap=200,
            )
        chunks = text_splitter.split_documents([doc])
        doc_chunks.extend(chunks)
    vectorstore = FAISS.from_documents(doc_chunks, embeddings)
    return vectorstore


def create_chain(vectorstore):
    """Original create_chain function"""
    llm = ChatGroq(
        model="llama-3.3-70b-versatile",
        temperature=0,
        max_tokens=4096,
    )
    retriever = vectorstore.as_retriever()
    memory = ConversationBufferMemory(
        llm=llm, output_key="answer", memory_key="chat_history", return_messages=True
    )
    chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=retriever,
        chain_type="map_reduce",
        memory=memory,
        verbose=True,
    )
    return chain


def create_chain_with_custom_retriever(vectorstore, k=20):
    """Create a chain with a custom retriever that returns more results"""
    llm = ChatGroq(
        model="llama-3.3-70b-versatile",
        temperature=0,
        max_tokens=4096,
    )
    # Create a custom retriever that returns more documents
    retriever = vectorstore.as_retriever(
        search_type="similarity",
        search_kwargs={"k": k},  # Return top k most similar chunks
    )
    memory = ConversationBufferMemory(
        llm=llm, output_key="answer", memory_key="chat_history", return_messages=True
    )
    chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=retriever,
        chain_type="map_reduce",  # Better for handling multiple documents
        memory=memory,
        verbose=True,
        return_source_documents=True,  # This helps with debugging
    )
    return chain


def handle_comprehensive_query(query: str, vectorstore) -> str:
    """Handle comprehensive queries with ALL CRM data for person searches"""
    # Keywords that indicate a comprehensive query
    comprehensive_keywords = [
        "list all",
        "all names",
        "all people",
        "everyone",
        "complete list",
        "full list",
        "entire",
        "total",
        "every",
        "tell me about",
        "phone number",
        "email",
        "contact",
    ]

    is_comprehensive = any(
        keyword in query.lower() for keyword in comprehensive_keywords
    )

    if is_comprehensive:
        all_docs = []

        # Check if this is a person/CRM query
        is_person_query = any(
            word in query.lower()
            for word in [
                "robert",
                "jason",
                "christina",
                "nicolas",
                "phone",
                "email",
                "contact",
                "tell me about",
            ]
        )

        if is_person_query:
            # For person queries, get ALL CRM chunks + some others
            if hasattr(vectorstore, "docstore") and hasattr(
                vectorstore, "index_to_docstore_id"
            ):
                # Get ALL CRM chunks
                for i in range(len(vectorstore.index_to_docstore_id)):
                    doc_id = vectorstore.index_to_docstore_id.get(i)
                    if doc_id and doc_id in vectorstore.docstore._dict:
                        doc = vectorstore.docstore._dict[doc_id]
                        if doc.metadata.get("file_type") == "xlsx":
                            all_docs.append(doc)

                # Also get some relevant non-CRM docs
                similarity_retriever = vectorstore.as_retriever(
                    search_type="similarity", search_kwargs={"k": 10}
                )
                other_docs = similarity_retriever.get_relevant_documents(query)
                for doc in other_docs:
                    if doc.metadata.get("file_type") != "xlsx":
                        all_docs.append(doc)
        else:
            # For non-person queries, use regular retrieval
            similarity_retriever = vectorstore.as_retriever(
                search_type="similarity", search_kwargs={"k": 40}
            )
            all_docs = similarity_retriever.get_relevant_documents(query)

        # Combine all content
        combined_content = "\n".join([doc.page_content for doc in all_docs])

        # Create a specialized prompt
        llm = ChatGroq(
            model="llama-3.3-70b-versatile",
            temperature=0,
            max_tokens=4096,
        )

        prompt = f"""
        Question: {query}
        
        Complete Database Content:
        {combined_content[:25000]}
        
        Answer the question based ONLY on the database content above. 
        For CRM records, look for: "FIRST_NAME: [name] | LAST_NAME: [name] | ... | PHONE1: [number]"
        
        Instructions:
        - Search through ALL the provided content carefully
        - For person queries, look through all the CRM records to find the exact person
        - For "list all" queries, extract ALL matching items from the content
        - Do not use any external knowledge - only use the information provided above
        - If you cannot find the information, say so clearly
        
        Answer:
        """

        try:
            response = llm.invoke(prompt)
            return response.content
        except Exception as e:
            return f"Error processing comprehensive query: {str(e)}"

    return None  # Not a comprehensive query


def initialize_storage():
    """Initialize storage and load existing data"""
    # Set up SQLite database
    setup_database()
    # Try to load existing vectorstore
    vectorstore = load_vectorstore(DATA_DIR)
    if vectorstore:
        st.session_state.vectorstore = vectorstore
        # Load URLs from database to session state
        documents, urls = load_processed_content()
        st.session_state.processed_urls = urls
        # Load chat history if available
        chat_history = load_chat_history()
        if chat_history:
            st.session_state.chat_history = chat_history
        st.sidebar.success(f"Loaded {len(urls)} previously processed items")
    return True


def get_content_statistics():
    """Get statistics about processed content"""
    conn = get_db_connection()
    cursor = conn.cursor()
    stats = {}
    # Count web pages
    cursor.execute("SELECT COUNT(*) FROM pages")
    stats["web_pages"] = cursor.fetchone()[0]
    # Count documents by type
    cursor.execute("SELECT file_type, COUNT(*) FROM documents GROUP BY file_type")
    doc_counts = cursor.fetchall()
    for file_type, count in doc_counts:
        stats[f"{file_type}_files"] = count
    return stats


def on_shutdown():
    close_db_connection()


def debug_retrieval(query: str, vectorstore):
    """Debug what chunks are being retrieved for a query"""
    st.write("## 🔍 Retrieval Debug")

    # Test the same retrieval as intelligent_rag_query
    retriever = vectorstore.as_retriever(
        search_type="mmr", search_kwargs={"k": 40, "fetch_k": 80, "lambda_mult": 0.3}
    )

    docs = retriever.get_relevant_documents(query)

    st.write(f"**Query**: {query}")
    st.write(f"**Retrieved {len(docs)} chunks:**")

    # Analyze what types of chunks we got
    chunk_types = {}
    for doc in docs:
        source = doc.metadata.get("source", "unknown")
        file_type = doc.metadata.get("file_type", "unknown")

        if "xlsx" in source.lower():
            chunk_type = "CRM"
        elif (
            "agreement" in source.lower()
            or doc.metadata.get("extraction_method") == "pymupdf_ocr"
        ):
            chunk_type = "Legal"
        elif "manual" in source.lower():
            chunk_type = "Manual"
        else:
            chunk_type = "Other"

        if chunk_type not in chunk_types:
            chunk_types[chunk_type] = 0
        chunk_types[chunk_type] += 1

    st.write("**Chunk breakdown:**")
    for chunk_type, count in chunk_types.items():
        st.write(f"- {chunk_type}: {count} chunks")

    # Show first few chunks
    st.write("**First 5 chunks:**")
    for i, doc in enumerate(docs[:5]):
        st.write(f"**Chunk {i + 1}:**")
        st.write(f"- Source: {doc.metadata.get('source', 'unknown')}")
        st.write(f"- Type: {doc.metadata.get('file_type', 'unknown')}")
        st.write(f"- Content preview: {doc.page_content[:100]}...")
        st.write("---")


def debug_crm_chunks(query: str, vectorstore):
    """Debug specifically what CRM records are being retrieved"""
    st.write("## 🔍 CRM Records Debug")

    retriever = vectorstore.as_retriever(
        search_type="mmr", search_kwargs={"k": 40, "fetch_k": 80, "lambda_mult": 0.3}
    )

    docs = retriever.get_relevant_documents(query)

    # Extract just the CRM chunks
    crm_chunks = [doc for doc in docs if doc.metadata.get("file_type") == "xlsx"]

    st.write(f"**Found {len(crm_chunks)} CRM chunks for query: '{query}'**")

    # Parse and show names from CRM chunks
    names_found = []
    for i, doc in enumerate(crm_chunks):
        content = doc.page_content

        # Extract names from this chunk
        import re

        name_matches = re.findall(
            r"FIRST_NAME:\s*([^|]+?)\s*\|\s*LAST_NAME:\s*([^|]+?)\s*\|", content
        )

        chunk_names = []
        for first, last in name_matches:
            full_name = f"{first.strip()} {last.strip()}"
            chunk_names.append(full_name)
            names_found.append(full_name)

        if i < 5:  # Show first 5 chunks
            st.write(f"**CRM Chunk {i + 1}** - Names: {chunk_names}")
            st.text_area(
                f"Content {i + 1}", content[:500], height=100, key=f"crm_debug_{i}"
            )

    # Check if Robert Earl specifically is in the retrieved names
    st.write("---")
    st.write(f"**All names found in retrieved CRM chunks:**")
    unique_names = sorted(set(names_found))
    for name in unique_names[:20]:  # Show first 20
        if "robert" in name.lower() and "earl" in name.lower():
            st.write(f"🎯 **{name}** ← TARGET FOUND!")
        elif "robert" in name.lower():
            st.write(f"👤 {name}")
        else:
            st.write(f"   {name}")

    if len(unique_names) > 20:
        st.write(f"... and {len(unique_names) - 20} more names")


############################################### END DEBUG  ############

# Register shutdown handler
import atexit

atexit.register(on_shutdown)


def handle_specific_person_query(query: str, vectorstore) -> str:
    """Handle queries about specific people - GUARANTEED WORKING VERSION"""
    try:
        # Extract the name from the query first
        import re

        # Simple name extraction
        words = query.split()
        potential_names = []

        # Look for pairs of capitalized words
        for i in range(len(words) - 1):
            if (
                words[i].replace(",", "").replace(".", "").isalpha()
                and words[i + 1].replace(",", "").replace(".", "").isalpha()
                and len(words[i]) > 1
                and len(words[i + 1]) > 1
            ):
                first = words[i].capitalize()
                last = words[i + 1].capitalize()

                # Skip obvious non-names
                if first.lower() not in [
                    "give",
                    "show",
                    "tell",
                    "what",
                    "have",
                    "all",
                    "the",
                    "information",
                ] and last.lower() not in [
                    "give",
                    "show",
                    "tell",
                    "what",
                    "have",
                    "all",
                    "the",
                    "information",
                ]:
                    potential_names.append(f"{first} {last}")

        # Also try single words that might be names
        for word in words:
            clean_word = word.replace(",", "").replace(".", "").capitalize()
            if (
                len(clean_word) > 2
                and clean_word.isalpha()
                and clean_word.lower()
                not in [
                    "give",
                    "show",
                    "tell",
                    "what",
                    "have",
                    "all",
                    "the",
                    "information",
                    "you",
                    "me",
                    "on",
                    "about",
                ]
            ):
                potential_names.append(clean_word)

        st.write(f"**Debug: Looking for these names: {potential_names}**")

        if not potential_names:
            return "I couldn't identify any names in your query. Please try asking like 'What is John Smith's phone number?'"

        # Get ALL the names from the database using our working function
        all_names = get_all_names_from_vectorstore_improved(vectorstore)

        st.write(f"**Debug: Total names in database: {len(all_names)}**")

        # Find exact matches (case insensitive)
        matched_names = []
        for search_name in potential_names:
            st.write(f"**Debug: Searching for '{search_name}'**")

        for db_name in all_names:
            # Exact match
            if search_name.lower() == db_name.lower():
                matched_names.append(db_name)
                st.write(f"**Debug: EXACT MATCH FOUND: {db_name}**")
            # First name match
            elif search_name.lower() == db_name.split()[0].lower():
                matched_names.append(db_name)
                st.write(f"**Debug: FIRST NAME MATCH: {db_name}**")
            # Last name match (check if search_name matches any part after first name)
            elif len(db_name.split()) > 1 and search_name.lower() in [
                part.lower() for part in db_name.split()[1:]
            ]:
                matched_names.append(db_name)
                st.write(f"**Debug: LAST NAME MATCH: {db_name}**")
            # Partial match (search_name is part of db_name)
            elif search_name.lower() in db_name.lower():
                matched_names.append(db_name)
                st.write(f"**Debug: PARTIAL MATCH FOUND: {db_name}**")
        # Remove duplicates
        matched_names = list(set(matched_names))

        if not matched_names:
            # Show similar names for debugging
            similar_names = [
                name
                for name in all_names
                if any(part.lower() in name.lower() for part in potential_names)
            ]
            return f"I couldn't find exact matches for: {potential_names}. Similar names in database: {similar_names[:10]}"

        st.write(f"**Debug: Found these matching names: {matched_names}**")

        # Now get the full records for these names
        all_content = []
        if hasattr(vectorstore, "docstore") and hasattr(
            vectorstore, "index_to_docstore_id"
        ):
            for i in range(len(vectorstore.index_to_docstore_id)):
                doc_id = vectorstore.index_to_docstore_id.get(i)
                if doc_id and doc_id in vectorstore.docstore._dict:
                    doc = vectorstore.docstore._dict[doc_id]
                    if doc.metadata.get("file_type") in ["xlsx", "xls"]:
                        all_content.append(doc.page_content)

        combined_content = "\n".join(all_content)

        # Find the full records for the matched names
        found_records = []
        for name in matched_names:
            # Split name into first and last
            name_parts = name.split()
            if len(name_parts) >= 2:
                first_name = name_parts[0]
                last_name = " ".join(name_parts[1:])  # Handle names like "Van Der Berg"

                # Search for the exact pattern in the content
                pattern = f"FIRST_NAME: {first_name} | LAST_NAME: {last_name}"

                # Find the record containing this pattern
                lines = combined_content.split("\n")
                for i, line in enumerate(lines):
                    if pattern in line:
                        # Get the full record (this line plus several after it)
                        record_lines = []
                        for j in range(i, min(i + 10, len(lines))):
                            record_lines.append(lines[j])
                            # Stop when we hit the next record or end
                            if j > i and (
                                "Record " in lines[j] and "FIRST_NAME:" in lines[j]
                            ):
                                break

                        found_records.append("\n".join(record_lines))
                        break

        if found_records:
            st.write(f"**Debug: Found {len(found_records)} complete records**")

            # Use LLM to format the response
            llm = ChatGroq(
                model="llama-3.3-70b-versatile",
                temperature=0,
                max_tokens=4096,
            )

            records_text = "\n\n".join(found_records)

            prompt = f"""
            Question: {query}
            
            Here are the complete database records for the person(s) you asked about:
            
            {records_text}
            
            Please provide a clear, organized response with all the information for the person requested.
            """

            response = llm.invoke(prompt)
            return response.content
        else:
            return f"Found matching names {matched_names} but couldn't retrieve their full records."

    except Exception as e:
        st.write(f"**Debug: Error: {str(e)}**")
        return f"Error: {str(e)}"


def search_manual_content(query: str, vectorstore) -> str:
    """Search specifically in manual content"""
    try:
        # Get only manual chunks
        manual_chunks = []
        if hasattr(vectorstore, "docstore") and hasattr(
            vectorstore, "index_to_docstore_id"
        ):
            for i in range(len(vectorstore.index_to_docstore_id)):
                doc_id = vectorstore.index_to_docstore_id.get(i)
                if doc_id and doc_id in vectorstore.docstore._dict:
                    doc = vectorstore.docstore._dict[doc_id]

                    # Only get manual chunks
                    if (
                        doc.metadata.get("file_type") == "pdf"
                        and "manual" in doc.metadata.get("source", "").lower()
                    ):
                        manual_chunks.append(doc.page_content)

        if not manual_chunks:
            return "I don't have any manual content available."

        # Combine manual content
        combined_manual = "\n".join(manual_chunks)

        # Use LLM with manual content only
        llm = ChatGroq(
            model="llama-3.3-70b-versatile",
            temperature=0,
            max_tokens=4096,
        )

        prompt = f"""
        Based on the following X570 AORUS ELITE motherboard manual content, please answer this question: {query}
        
        Manual Content:
        {combined_manual[:15000]}  # Limit to avoid token limits
        
        Please provide a detailed answer based on the manual content above.
        """

        response = llm.invoke(prompt)
        return response.content

    except Exception as e:
        return f"Error searching manual: {str(e)}"


def process_pdf_with_ocr(file_path: str) -> Document:
    """Process a PDF using OCR if regular text extraction fails"""
    try:
        import pytesseract
        from pdf2image import convert_from_path
        from PIL import Image
        import tempfile
        import os

        st.write(f"**Processing PDF with OCR**: {os.path.basename(file_path)}")

        # First try regular Docling extraction
        try:
            converter = DocumentConverter()
            result = converter.convert(file_path)

            if hasattr(result.document, "text"):
                text_content = result.document.text
            elif hasattr(result.document, "get_text"):
                text_content = result.document.get_text()
            else:
                text_content = str(result.document)

            # Check if we got meaningful text (more than just image tags)
            if len(text_content) > 500 and "<!-- image -->" not in text_content:
                st.write("✅ **Regular text extraction successful**")
                return Document(
                    page_content=text_content,
                    metadata={
                        "source": file_path,
                        "file_name": os.path.basename(file_path),
                        "file_type": "pdf",
                        "extraction_method": "docling",
                    },
                )
        except Exception as e:
            st.write(f"⚠️ **Regular extraction failed**: {str(e)}")

        st.write("🔍 **Attempting OCR extraction...**")

        # Convert PDF pages to images
        with tempfile.TemporaryDirectory() as temp_dir:
            try:
                # Convert PDF to images
                pages = convert_from_path(file_path, dpi=300)
                st.write(f"📄 **Converted to {len(pages)} images**")

                extracted_text = []

                for i, page in enumerate(pages):
                    st.write(f"🔍 **Processing page {i + 1}/{len(pages)}...**")

                    # Save page as temporary image
                    temp_image_path = os.path.join(temp_dir, f"page_{i + 1}.png")
                    page.save(temp_image_path, "PNG")

                    # Extract text using OCR
                    try:
                        page_text = pytesseract.image_to_string(page)
                        if page_text.strip():
                            extracted_text.append(f"\n--- Page {i + 1} ---\n")
                            extracted_text.append(page_text.strip())
                        else:
                            st.write(f"⚠️ **No text found on page {i + 1}**")
                    except Exception as e:
                        st.write(f"❌ **OCR failed for page {i + 1}**: {str(e)}")
                        continue

                if extracted_text:
                    full_text = "\n".join(extracted_text)
                    st.write(
                        f"✅ **OCR successful! Extracted {len(full_text)} characters**"
                    )

                    # Show preview
                    st.write("**OCR Text Preview (first 500 characters):**")
                    st.text_area("OCR Preview", full_text[:500], height=150)

                    return Document(
                        page_content=full_text,
                        metadata={
                            "source": file_path,
                            "file_name": os.path.basename(file_path),
                            "file_type": "pdf",
                            "extraction_method": "ocr",
                            "pages_processed": len(pages),
                            "text_length": len(full_text),
                        },
                    )
                else:
                    st.error("❌ **No text could be extracted via OCR**")
                    return None

            except Exception as e:
                st.error(f"❌ **PDF to image conversion failed**: {str(e)}")
                return None

    except ImportError:
        st.error(
            "❌ **OCR libraries not installed**. Please install: pip install pytesseract pillow pdf2image"
        )
        return None
    except Exception as e:
        st.error(f"❌ **OCR processing failed**: {str(e)}")
        return None


# Streamlit app configuration
st.set_page_config(page_title="Chat with Website", page_icon="", layout="wide")
st.title("CLAUDIA  🦙(LLAMA 3.3)")

# Initialize session state
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "processed_urls" not in st.session_state:
    st.session_state.processed_urls = []

# Initialize storage when app starts
if "storage_initialized" not in st.session_state:
    initialize_storage()
    st.session_state.storage_initialized = True

# Document loading options
st.subheader("Load Content")

# Sitemap input
sitemap_col, sitemap_options_col = st.columns([3, 1])
with sitemap_col:
    sitemap_url = st.text_input(
        "Enter website sitemap URL:",
        placeholder="e.g., https://example.com/sitemap.xml",
    )
with sitemap_options_col:
    max_urls = st.slider(
        "Max URLs",
        min_value=1,
        max_value=100,
        value=50,
        help="Limit the number of URLs to process",
    )

# Local HTML directory input
html_col, html_options_col = st.columns([3, 1])
with html_col:
    html_dir = st.text_input(
        "Enter path to HTML files directory:",
        placeholder="e.g., C:/Documents/html_files or /home/user/html_files",
    )
with html_options_col:
    max_html_files = st.slider(
        "Max HTML files",
        min_value=1,
        max_value=100,
        value=50,
        help="Limit the number of HTML files to process",
    )

# Local DOCX directory input
docx_col, docx_options_col = st.columns([3, 1])
with docx_col:
    docx_dir = st.text_input(
        "Enter path to DOCX files directory:",
        placeholder="e.g., C:/Documents/docx_files or /home/user/docx_files",
    )
with docx_options_col:
    max_docx_files = st.slider(
        "Max DOCX files",
        min_value=1,
        max_value=100,
        value=50,
        help="Limit the number of DOCX files to process",
    )

# Local PDF directory input
pdf_col, pdf_options_col = st.columns([3, 1])
with pdf_col:
    pdf_dir = st.text_input(
        "Enter path to PDF files directory:",
        placeholder="e.g., C:/Documents/pdf_files or /home/user/pdf_files",
    )
with pdf_options_col:
    max_pdf_files = st.slider(
        "Max PDF files",
        min_value=1,
        max_value=100,
        value=50,
        help="Limit the number of PDF files to process",
    )

# Local PowerPoint directory input
pptx_col, pptx_options_col = st.columns([3, 1])
with pptx_col:
    pptx_dir = st.text_input(
        "Enter path to PowerPoint files directory:",
        placeholder="e.g., C:/Documents/pptx_files or /home/user/pptx_files",
    )
with pptx_options_col:
    max_pptx_files = st.slider(
        "Max PPTX files",
        min_value=1,
        max_value=100,
        value=50,
        help="Limit the number of PowerPoint files to process",
    )

# Local Excel directory input
xlsx_col, xlsx_options_col = st.columns([3, 1])
with xlsx_col:
    xlsx_dir = st.text_input(
        "Enter path to Excel files directory:",
        placeholder="e.g., C:/Documents/xlsx_files or /home/user/xlsx_files",
    )
with xlsx_options_col:
    max_xlsx_files = st.slider(
        "Max Excel files",
        min_value=1,
        max_value=100,
        value=50,
        help="Limit the number of Excel files to process",
    )

# Process buttons
col1, col2, col3, col4, col5, col6 = st.columns(6)
with col1:
    process_sitemap_button = st.button(
        "Process Sitemap", type="primary", use_container_width=True
    )
with col2:
    process_html_button = st.button(
        "Process HTML", type="primary", use_container_width=True
    )
with col3:
    process_docx_button = st.button(
        "Process DOCX", type="primary", use_container_width=True
    )
with col4:
    process_pdf_button = st.button(
        "Process PDF", type="primary", use_container_width=True
    )
with col5:
    process_pptx_button = st.button(
        "Process PowerPoint", type="primary", use_container_width=True
    )
with col6:
    process_xlsx_button = st.button(
        "Process Excel", type="primary", use_container_width=True
    )

# with col4:
#     # Add this test button for OCR Programs Installed
#     if st.button("🧪 Test Poppler", key="test_poppler"):
#         try:
#             import subprocess

#             # Test if poppler is accessible
#             result = subprocess.run(["pdftoppm", "-h"], capture_output=True, text=True)
#             if result.returncode == 0:
#                 st.success("✅ Poppler is installed and accessible!")
#             else:
#                 st.error("❌ Poppler not found in PATH")
#         except FileNotFoundError:
#             st.error("❌ Poppler not installed or not in PATH")
#             st.write("**Install instructions:**")
#             st.code("sudo apt install poppler-utils")
#         except Exception as e:
#             st.error(f"❌ Error testing Poppler: {e}")

# with col4:
#     if st.button("🧪 Test OCR", key="test_ocr_setup"):
#         try:
#             import pytesseract
#             import pdf2image
#             from PIL import Image

#             st.success("✅ All OCR libraries are installed!")
#             # Test Tesseract
#             version = pytesseract.get_tesseract_version()
#             st.write(f"📝 Tesseract version: {version}")
#         except ImportError as e:
#             st.error(f"❌ Missing library: {e}")
#             st.write("Install with: `pip install pytesseract pillow pdf2image`")
#         except Exception as e:
#             st.error(f"❌ Tesseract not found: {e}")
#             st.write("Install Tesseract: https://github.com/UB-Mannheim/tesseract/wiki")

# Process sitemap
if process_sitemap_button and sitemap_url:
    try:
        documents, processed_urls = process_sitemap(sitemap_url, max_urls)
        if documents:
            st.success(
                f"Successfully processed {len(documents)} pages from {len(processed_urls)} URLs"
            )
            # Create or update vectorstore
            if "vectorstore" not in st.session_state:
                st.session_state.vectorstore = setup_vectorstore(documents)
            else:
                st.session_state.vectorstore.add_documents(documents)
            # Save data to persistent storage
            save_vectorstore(st.session_state.vectorstore)
            save_processed_content(sitemap_url, processed_urls, documents)
            st.session_state.processed_urls.extend(processed_urls)
            # Show processed URLs
            with st.expander("Processed URLs", expanded=True):
                for url in processed_urls:
                    st.write(url)
    except Exception as e:
        st.error(f"Error processing sitemap: {str(e)}")

# Process HTML files
if process_html_button and html_dir:
    try:
        documents, processed_files = process_local_directory(
            html_dir, "html", max_html_files
        )
        if documents:
            st.success(f"Successfully processed {len(documents)} HTML files")
            # Create or update vectorstore
            if "vectorstore" not in st.session_state:
                st.session_state.vectorstore = setup_vectorstore(documents)
            else:
                st.session_state.vectorstore.add_documents(documents)
            # Save data to persistent storage
            save_vectorstore(st.session_state.vectorstore)
            save_local_documents(documents, processed_files, "html")
            st.session_state.processed_urls.extend(processed_files)
            # Show processed files
            with st.expander("Processed HTML Files", expanded=True):
                for file_path in processed_files:
                    st.write(os.path.basename(file_path))
    except Exception as e:
        st.error(f"Error processing HTML files: {str(e)}")

# Process DOCX files
if process_docx_button and docx_dir:
    try:
        documents, processed_files = process_local_directory(
            docx_dir, "docx", max_docx_files
        )
        if documents:
            st.success(f"Successfully processed {len(documents)} DOCX files")
            # Create or update vectorstore
            if "vectorstore" not in st.session_state:
                st.session_state.vectorstore = setup_vectorstore(documents)
            else:
                st.session_state.vectorstore.add_documents(documents)
            # Save data to persistent storage
            save_vectorstore(st.session_state.vectorstore)
            save_local_documents(documents, processed_files, "docx")
            st.session_state.processed_urls.extend(processed_files)
            # Show processed files
            with st.expander("Processed DOCX Files", expanded=True):
                for file_path in processed_files:
                    st.write(os.path.basename(file_path))
    except Exception as e:
        st.error(f"Error processing DOCX files: {str(e)}")

# Process PDF files
if process_pdf_button and pdf_dir:
    try:
        documents, processed_files = process_local_directory(
            pdf_dir, "pdf", max_pdf_files
        )
        if documents:
            st.success(f"Successfully processed {len(documents)} PDF files")
            # Create or update vectorstore
            if "vectorstore" not in st.session_state:
                st.session_state.vectorstore = setup_vectorstore(documents)
            else:
                st.session_state.vectorstore.add_documents(documents)
            # Save data to persistent storage
            save_vectorstore(st.session_state.vectorstore)
            save_local_documents(documents, processed_files, "pdf")
            st.session_state.processed_urls.extend(processed_files)
            # Show processed files
            with st.expander("Processed PDF Files", expanded=True):
                for file_path in processed_files:
                    st.write(os.path.basename(file_path))
    except Exception as e:
        st.error(f"Error processing PDF files: {str(e)}")

# Process PowerPoint files
if process_pptx_button and pptx_dir:
    try:
        documents, processed_files = process_local_directory(
            pptx_dir, "pptx", max_pptx_files
        )
        if documents:
            st.success(f"Successfully processed {len(documents)} PowerPoint files")
            # Create or update vectorstore
            if "vectorstore" not in st.session_state:
                st.session_state.vectorstore = setup_vectorstore(documents)
            else:
                st.session_state.vectorstore.add_documents(documents)
            # Save data to persistent storage
            save_vectorstore(st.session_state.vectorstore)
            save_local_documents(documents, processed_files, "pptx")
            st.session_state.processed_urls.extend(processed_files)
            # Show processed files
            with st.expander("Processed PowerPoint Files", expanded=True):
                for file_path in processed_files:
                    st.write(os.path.basename(file_path))
    except Exception as e:
        st.error(f"Error processing PowerPoint files: {str(e)}")

# Process Excel files
if process_xlsx_button and xlsx_dir:
    try:
        documents, processed_files = process_local_directory(
            xlsx_dir, "xlsx", max_xlsx_files
        )
        if documents:
            st.success(f"Successfully processed {len(documents)} Excel files")
            # Create or update vectorstore
            if "vectorstore" not in st.session_state:
                st.session_state.vectorstore = setup_vectorstore(documents)
            else:
                st.session_state.vectorstore.add_documents(documents)
            # Save data to persistent storage
            save_vectorstore(st.session_state.vectorstore)
            save_local_documents(documents, processed_files, "xlsx")
            st.session_state.processed_urls.extend(processed_files)
            # Show processed files
            with st.expander("Processed Excel Files", expanded=True):
                for file_path in processed_files:
                    st.write(os.path.basename(file_path))
    except Exception as e:
        st.error(f"Error processing Excel files: {str(e)}")


# # Add this as a completely separate section - put it after your other diagnostics
# st.write("---")
# st.subheader("🔍 OCR PDF Processing")

# # Simple file path input
# ocr_pdf_path = st.text_input(
#     "PDF file path for OCR:",
#     value="/mnt/c/AI/CRMRecords/agreement-S.pdf",
#     key="simple_ocr_path",
# )

# # Simple process button
# if st.button("Start OCR Processing", key="simple_ocr_start"):
#     if ocr_pdf_path and os.path.exists(ocr_pdf_path):
#         st.write(f"🔍 **Starting OCR on**: {os.path.basename(ocr_pdf_path)}")

#         try:
#             # Show that we're starting
#             with st.status("Processing PDF with OCR...", expanded=True) as status:
#                 st.write("📄 Converting PDF to images...")

#                 # Import here to check for errors
#                 try:
#                     import pytesseract
#                     import pdf2image
#                     from PIL import Image

#                     st.write("✅ OCR libraries loaded")
#                 except ImportError as e:
#                     st.error(f"❌ Missing libraries: {e}")
#                     st.stop()

#                 # Simple OCR processing
#                 st.write("🔄 Starting OCR extraction...")

#                 # Convert PDF to images
#                 pages = pdf2image.convert_from_path(ocr_pdf_path, dpi=200)
#                 st.write(f"📄 Converted to {len(pages)} page images")

#                 # Extract text from each page
#                 all_text = []
#                 for i, page in enumerate(pages):
#                     st.write(f"🔍 Processing page {i + 1}/{len(pages)}")
#                     try:
#                         page_text = pytesseract.image_to_string(page)
#                         if page_text.strip():
#                             all_text.append(
#                                 f"\n--- Page {i + 1} ---\n{page_text.strip()}"
#                             )
#                         else:
#                             st.write(f"⚠️ No text on page {i + 1}")
#                     except Exception as e:
#                         st.write(f"❌ OCR failed on page {i + 1}: {e}")

#                 if all_text:
#                     combined_text = "\n".join(all_text)
#                     st.write(
#                         f"✅ **OCR Complete!** Extracted {len(combined_text)} characters"
#                     )

#                     # Store result
#                     st.session_state["ocr_text"] = combined_text
#                     st.session_state["ocr_file"] = ocr_pdf_path

#                     status.update(label="✅ OCR Processing Complete!", state="complete")
#                 else:
#                     st.error("❌ No text could be extracted")

#         except Exception as e:
#             st.error(f"❌ **OCR Error**: {str(e)}")
#             import traceback

#             st.code(traceback.format_exc())
#     else:
#         st.error("❌ PDF file not found!")

# # Show results if available
# if "ocr_text" in st.session_state and st.session_state["ocr_text"]:
#     st.write("---")
#     st.write("## 📄 OCR Results")

#     ocr_text = st.session_state["ocr_text"]
#     ocr_file = st.session_state["ocr_file"]

#     st.write(f"**File**: {os.path.basename(ocr_file)}")
#     st.write(f"**Extracted**: {len(ocr_text)} characters")

#     # Show preview
#     st.text_area(
#         "OCR Text Preview", ocr_text[:1000], height=200, key="ocr_result_preview"
#     )

#     # Add to vector store button
#     if st.button("Add OCR Text to Vector Store", key="add_ocr_final"):
#         try:
#             # Create document
#             doc = Document(
#                 page_content=ocr_text,
#                 metadata={
#                     "source": ocr_file,
#                     "file_name": os.path.basename(ocr_file),
#                     "file_type": "pdf",
#                     "extraction_method": "ocr",
#                     "text_length": len(ocr_text),
#                 },
#             )

#             # Create chunks
#             text_splitter = CharacterTextSplitter(
#                 separator="\n",
#                 chunk_size=1000,
#                 chunk_overlap=200,
#             )
#             chunks = text_splitter.split_documents([doc])

#             # Add to vector store
#             if "vectorstore" in st.session_state:
#                 st.session_state.vectorstore.add_documents(chunks)
#                 save_vectorstore(st.session_state.vectorstore)
#                 save_local_documents([doc], [ocr_file], "pdf")

#                 st.success(
#                     f"✅ **Success!** Added {len(chunks)} chunks to vector store"
#                 )

#                 st.success(f"✅ **Added {len(chunks)} chunks to vector store!**")

#                 # Clear results without immediate rerun
#                 del st.session_state["ocr_text"]
#                 del st.session_state["ocr_file"]

#                 # Show instruction instead of auto-rerun
#                 st.info("🔄 **Refresh the page to clear this section**")
#             else:
#                 st.error("❌ No vector store found!")

#         except Exception as e:
#             st.error(f"❌ Error adding to vector store: {str(e)}")

#     # Clear results button
#     if st.button("Clear OCR Results", key="clear_ocr_simple"):
#         del st.session_state["ocr_text"]
#         del st.session_state["ocr_file"]
#         st.rerun()


# Show loaded documents/URLs in sidebar
with st.sidebar:
    st.subheader("Processed Content")
    if st.session_state.processed_urls:
        # Get content statistics
        stats = get_content_statistics()
        st.write(f"**Total items:** {len(st.session_state.processed_urls)}")
        # Show breakdown by type
        if stats.get("web_pages", 0) > 0:
            st.write(f"📄 Web pages: {stats['web_pages']}")
        if stats.get("html_files", 0) > 0:
            st.write(f"📝 HTML files: {stats['html_files']}")
        if stats.get("docx_files", 0) > 0:
            st.write(f"📘 DOCX files: {stats['docx_files']}")
        if stats.get("pdf_files", 0) > 0:
            st.write(f"📕 PDF files: {stats['pdf_files']}")
        if stats.get("pptx_files", 0) > 0:
            st.write(f"📊 PowerPoint files: {stats['pptx_files']}")
        if stats.get("xlsx_files", 0) > 0:
            st.write(f"📈 Excel files: {stats['xlsx_files']}")
        with st.expander("View all items", expanded=False):
            for item in st.session_state.processed_urls:
                if item.startswith("http"):
                    st.write(f"🌐 {item}")
                else:
                    st.write(f"📁 {os.path.basename(item)}")

        # Initialize confirmation state
        if "db_delete_confirmation" not in st.session_state:
            st.session_state.db_delete_confirmation = False

        col1, col2 = st.columns(2)
        with col1:
            # Display different buttons based on confirmation state
            if not st.session_state.db_delete_confirmation:
                # Initial button
                if st.button("Clear Database", type="primary"):
                    st.session_state.db_delete_confirmation = True
                    st.rerun()
            else:
                # Show warning and confirmation buttons
                st.warning(
                    "⚠️ Are you sure you want to delete all data? This cannot be undone."
                )
                # Yes button
                if st.button("Yes, Delete", type="primary"):
                    # Clear tables but keep structure
                    conn = get_db_connection()
                    cursor = conn.cursor()
                    # Delete in the correct order to respect foreign key constraints
                    cursor.execute("DELETE FROM pages")
                    cursor.execute("DELETE FROM websites")
                    cursor.execute("DELETE FROM documents")
                    conn.commit()
                    # Remove vector store files
                    import shutil

                    if os.path.exists(os.path.join(DATA_DIR, "faiss_index")):
                        shutil.rmtree(os.path.join(DATA_DIR, "faiss_index"))
                    # Clear session state but keep chat history
                    for key in [
                        "vectorstore",
                        "conversation_chain",
                        "processed_urls",
                        "db_delete_confirmation",
                    ]:
                        if key in st.session_state:
                            del st.session_state[key]
                    st.success("Database cleared successfully!")
                    st.rerun()
                # No button
                if st.button("No, Cancel"):
                    st.session_state.db_delete_confirmation = False
                    st.rerun()

        with col2:
            if st.button("Clear Chat History", type="secondary"):
                # Clear only chat history
                conn = get_db_connection()
                cursor = conn.cursor()
                cursor.execute("DELETE FROM chat_history")
                conn.commit()
                # Clear chat history from session state
                st.session_state.chat_history = []
                # If conversation chain exists, reset its memory
                if "conversation_chain" in st.session_state:
                    # Create a new memory object
                    llm = ChatGroq(model="llama-3.3-70b-versatile", temperature=0)
                    memory = ConversationBufferMemory(
                        llm=llm,
                        output_key="answer",
                        memory_key="chat_history",
                        return_messages=True,
                    )
                    # Update the chain with the new memory
                    st.session_state.conversation_chain.memory = memory
                st.success("Chat history cleared!")
                st.rerun()
    else:
        st.info("No content processed yet. Enter paths or URLs above to begin.")


# Add this after the sidebar section and before the chat interface
if "vectorstore" in st.session_state:
    st.divider()
    export_container = st.container()
    with export_container:
        st.subheader("Data Export")
        # Create two columns for export buttons
        col1, col2 = st.columns(2)
        # Export Vector Database
        with col1:
            # Add export button for FAISS vectors
            if st.button("Export Vectors to CSV", type="primary"):
                with st.spinner("Exporting vectors..."):
                    # Define the export function
                    def export_faiss_vectors(filename="faiss_vectors.csv"):
                        """Export FAISS vectors to CSV for external visualization"""
                        try:
                            vs = st.session_state.vectorstore
                            # Get all the vectors and their IDs
                            vectors = []
                            ids = []
                            metadata = []
                            # Different vectorstores have different structures
                            if hasattr(vs, "index_to_docstore_id") and hasattr(
                                vs, "index"
                            ):
                                for i in range(len(vs.index_to_docstore_id)):
                                    doc_id = vs.index_to_docstore_id.get(i)
                                    if doc_id and doc_id in vs.docstore._dict:
                                        # Get the document
                                        doc = vs.docstore._dict[doc_id]
                                        # Get the vector
                                        vector = vs.index.reconstruct(i)
                                        vectors.append(vector)
                                        ids.append(doc_id)
                                        # Enhanced metadata for different document types
                                        meta = {
                                            "source": doc.metadata.get("source", ""),
                                            "content_preview": doc.page_content[:100],
                                        }
                                        if doc.metadata.get("type") == "web":
                                            meta["title"] = doc.metadata.get(
                                                "title", ""
                                            )
                                            meta["url"] = doc.metadata.get("url", "")
                                        else:
                                            meta["file_name"] = doc.metadata.get(
                                                "file_name", ""
                                            )
                                            meta["file_type"] = doc.metadata.get(
                                                "file_type", ""
                                            )
                                        metadata.append(meta)
                            # Export to CSV
                            import pandas as pd
                            import numpy as np

                            # Create a dataframe with metadata
                            meta_df = pd.DataFrame(metadata)
                            # Create a dataframe with vectors
                            vector_df = pd.DataFrame(np.array(vectors))
                            # Combine the dataframes
                            result_df = pd.concat([meta_df, vector_df], axis=1)
                            # Save to CSV
                            result_df.to_csv(filename, index=False)
                            return (
                                True,
                                f"Exported {len(vectors)} vectors to {filename}",
                            )
                        except Exception as e:
                            return False, f"Error exporting vectors: {str(e)}"

                    # Export the vectors to a file in the data directory
                    export_path = os.path.join(DATA_DIR, "faiss_vectors.csv")
                    success, message = export_faiss_vectors(export_path)
                    if success:
                        st.success(message)
                        # Create a download link for the exported file
                        with open(export_path, "rb") as file:
                            st.download_button(
                                label="Download Vector CSV",
                                data=file,
                                file_name="faiss_vectors.csv",
                                mime="text/csv",
                            )
                    else:
                        st.error(message)

        # Export SQLite Database
        with col2:
            if st.button("Export SQLite Database", type="primary"):
                with st.spinner("Preparing SQLite database for export..."):
                    try:
                        # Create a copy of the database for export
                        import shutil

                        export_db_path = os.path.join(
                            DATA_DIR, "website_data_export.db"
                        )
                        # Close any existing connections before copying
                        close_db_connection()
                        # Copy the database file
                        shutil.copy2(DB_PATH, export_db_path)
                        st.success("SQLite database prepared for download!")
                        # Create download button for the SQLite file
                        with open(export_db_path, "rb") as file:
                            st.download_button(
                                label="Download SQLite Database",
                                data=file,
                                file_name="website_data.db",
                                mime="application/x-sqlite3",
                            )
                    except Exception as e:
                        st.error(f"Error exporting SQLite database: {str(e)}")


# Initialize conversation chain if vectorstore exists
if "vectorstore" in st.session_state and "conversation_chain" not in st.session_state:
    st.session_state.conversation_chain = create_chain_with_custom_retriever(
        st.session_state.vectorstore,
        k=50,  # Retrieve top 50 chunks
    )

###################### Chat Handler (Begin) ####################


# Chat interface
if "vectorstore" in st.session_state:
    st.divider()
    st.subheader("Ask Claudia - Your AI Assistant")

    # Display chat history
    for message in st.session_state.chat_history:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    # Chat input
user_input = st.chat_input("Ask about the website content...")

if user_input and user_input.strip():
    # Add user message to chat history
    st.session_state.chat_history.append({"role": "user", "content": user_input})
    with st.chat_message("user"):
        st.write("Lowly Human")
        st.markdown(user_input)

    # Generate assistant response
    with st.chat_message("assistant"):
        st.write("Claudia")
        with st.spinner("Thinking..."):
            user_query_lower = user_input.lower()

            # VERY SPECIFIC routing for name-related queries only
            if any(
                keyword in user_query_lower
                for keyword in [
                    "list all names",
                    "list the names",
                    "all names",
                    "show all names",
                ]
            ):
                # Handle "list all names" requests
                names = get_all_names_from_vectorstore_improved(
                    st.session_state.vectorstore
                )
                if names:
                    assistant_response = (
                        f"Here are all {len(names)} names in the database:\n\n"
                    )
                    for i, name in enumerate(names, 1):
                        assistant_response += f"{i}. {name}\n"
                else:
                    assistant_response = "I couldn't find any names in the database."
                st.markdown(assistant_response)

            elif any(
                keyword in user_query_lower
                for keyword in ["how many names", "count names", "number of names"]
            ):
                # Handle "how many names" requests
                names = get_all_names_from_vectorstore_improved(
                    st.session_state.vectorstore
                )
                assistant_response = f"There are {len(names)} names in the database."
                st.markdown(assistant_response)

            else:
                # Use the original simple approach for everything else
                comprehensive_response = handle_comprehensive_query(
                    user_input, st.session_state.vectorstore
                )

                if comprehensive_response:
                    assistant_response = comprehensive_response
                    st.markdown(assistant_response)
                else:
                    # Use regular chain
                    response = st.session_state.conversation_chain(
                        {"question": user_input}
                    )
                    assistant_response = response["answer"]
                    st.markdown(assistant_response)

                    # Show source documents for debugging
                    if "source_documents" in response:
                        with st.expander("Source Documents Used", expanded=False):
                            for i, doc in enumerate(response["source_documents"]):
                                st.write(f"**Source {i + 1}:**")
                                st.write(
                                    doc.page_content[:300] + "..."
                                    if len(doc.page_content) > 300
                                    else doc.page_content
                                )

            # Add assistant response to chat history
            st.session_state.chat_history.append(
                {"role": "assistant", "content": assistant_response}
            )

            # Save chat history to database
            save_chat_history(st.session_state.chat_history)

##################################### Chat Handler (End)
